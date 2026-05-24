#!/usr/bin/env python3
"""
PPT Storyboard Generator for UHN Accessibility First Course Series
Converts Master Storyboard markdown into branded PowerPoint for Storyline 360 import.

Each screen becomes one or more slides:
- Content screens → 1 slide
- Question screens → question slide + correct feedback slide + incorrect feedback slide
- Branching scenarios → setup slide + consequence slides per choice (A/B/C) + debrief slide

Usage:
    python generate_ppt.py --guide 1
    python generate_ppt.py --guide 1 --dry-run    # Preview slide plan without generating

Requirements:
    pip install python-pptx Pillow
"""

import argparse
import re
import sys
from pathlib import Path
from copy import deepcopy

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Error: 'python-pptx' package not installed. Run: pip install python-pptx")
    sys.exit(1)

BASE_DIR = Path(__file__).resolve().parent.parent.parent
BUILD_OUTPUT = BASE_DIR / "05-build-output"
TEMPLATE_PATH = BASE_DIR / "02-branding-and-style" / "uhn-templates" / "UHN_Presentation_Template_Primary (1).pptx"
LOGO_PATH = BASE_DIR / "02-branding-and-style" / "logos" / "uhn_logo.png"

GUIDE_FOLDERS = {
    1: "01-Foundations-of-Disability-Inclusion-and-Accessible-Design",
    2: "02-Perceptions-Attitudes-and-Barriers",
    3: "03-Vision-Disabilities",
    4: "04-Sensory-Hearing-and-Communication-Disabilities",
    5: "05-Physical-Disabilities-and-Mobility",
    6: "06-Mental-Health-Disabilities",
    7: "07-Intellectual-Developmental-and-Learning-Disabilities",
    8: "08-Non-Visible-Disabilities",
    9: "09-Aging-Disability-and-Intersectionality",
    10: "10-Engaging-with-Confidence-and-Respect",
    11: "11-Service-Animals-Guide-Dogs-and-Non-Service-Animals",
    12: "12-Support-Persons",
    13: "13-Assistive-Devices",
    14: "14-Communication-and-Information-Accessibility",
    15: "15-Neurodiversity-and-Sensory-Regulation",
    16: "16-Trauma-Informed-Accessibility",
    17: "17-Accessibility-in-Crisis-Situations-and-De-escalation",
    18: "18-Indigenous-Peoples-and-Accessibility",
}

# UHN Brand Colours
UHN_NAVY = RGBColor(0x19, 0x28, 0x58)
UHN_RED = RGBColor(0xC0, 0x23, 0x3B)
UHN_COBALT = RGBColor(0x24, 0x5B, 0xAA)
UHN_TEAL = RGBColor(0x00, 0xA5, 0xA8)
UHN_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
UHN_LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
CORRECT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
INCORRECT_RED = RGBColor(0xC6, 0x28, 0x28)

# Layout indices from UHN template
LAYOUT_TITLE = 1          # Title Slide - with graphic and Image-3
LAYOUT_SECTION = 7        # Section Title - UHN Blue
LAYOUT_TEXT_BLUE = 9      # TextOnly slide-Blue (title + body)
LAYOUT_GRAPHIC_BLUE = 10  # Graphic with text slide-Blue (two columns)
LAYOUT_TEXT_NEUTRAL = 21  # TextOnly slide-Neutral
LAYOUT_IMAGE_SLIDE = 29   # TextOnly slide-Image in Graphic
LAYOUT_CLOSING = 39       # Closing Slide


def set_text_props(run, font_name="Arial", font_size=14, bold=False, italic=False, color=None):
    """Apply consistent text formatting to a run."""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def add_text_to_placeholder(placeholder, text, font_size=14, bold=False, color=None, alignment=None):
    """Set text on a placeholder with formatting."""
    placeholder.text = ""
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_text_props(run, font_size=font_size, bold=bold, color=color)
    return p


def add_bullet_list(placeholder, items, font_size=12, color=None, bold_first=False):
    """Add a bulleted list to a placeholder."""
    tf = placeholder.text_frame
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)

        if bold_first and ":" in item:
            label, rest = item.split(":", 1)
            run_bold = p.add_run()
            run_bold.text = label + ":"
            set_text_props(run_bold, font_size=font_size, bold=True, color=color)
            run_rest = p.add_run()
            run_rest.text = rest
            set_text_props(run_rest, font_size=font_size, color=color)
        else:
            run = p.add_run()
            run.text = item
            set_text_props(run, font_size=font_size, color=color)


def add_multiline_text(placeholder, lines, font_size=12, color=None):
    """Add multiple paragraphs to a placeholder."""
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        set_text_props(run, font_size=font_size, color=color)


# ---------------------------------------------------------------------------
# Storyboard parser
# ---------------------------------------------------------------------------

def parse_storyboard(filepath: Path) -> dict:
    """Parse the master storyboard markdown into structured data."""
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    # Extract header info
    info = {}
    title_match = re.search(r'^# (.+)', content, re.MULTILINE)
    if title_match:
        info['title'] = title_match.group(1).strip()

    seat_match = re.search(r'Seat Time\s*\|\s*(.+?)(?:\s*\|)', content)
    if seat_match:
        info['seat_time'] = seat_match.group(1).strip()

    screens_match = re.search(r'Screens\s*\|\s*(\d+)', content)
    if screens_match:
        info['total_screens'] = int(screens_match.group(1))

    clo_match = re.search(r'\*\*CLOs?:\*\*\s*(.+?)$', content, re.MULTILINE)
    if clo_match:
        info['clos'] = clo_match.group(1).strip()

    # Parse individual screens
    screen_sections = re.split(r'(?=^## Screen \d)', content, flags=re.MULTILINE)
    screens = []

    for section in screen_sections:
        if not section.strip().startswith('## Screen'):
            continue

        screen = parse_screen(section)
        if screen:
            screens.append(screen)

    return {'info': info, 'screens': screens}


def parse_screen(section: str) -> dict:
    """Parse a single screen section into structured data."""
    screen = {}

    # Screen number and title
    header_match = re.search(r'## Screen ([\d.]+)\s*[—–-]\s*(.+)', section)
    if header_match:
        screen['number'] = header_match.group(1).strip()
        screen['title'] = header_match.group(2).strip()

    # Determine screen type
    screen['type'] = 'content'  # default
    if 'Scenario — Branching' in section or 'Branching' in section:
        screen['type'] = 'scenario'
    elif 'Assessment' in section or 'Knowledge Check' in section or 'Question' in section.split('\n')[0]:
        screen['type'] = 'question'
    elif 'Reflection' in section:
        screen['type'] = 'reflection'
    elif 'MAP' in section or 'Action Planning' in section:
        screen['type'] = 'action_planning'
    elif 'Welcome' in section:
        screen['type'] = 'welcome'
    elif 'Learning Objectives' in section or 'Objectives' in section.split('\n')[0]:
        screen['type'] = 'objectives'
    elif 'Completion' in section or 'Resources' in section.split('\n')[0]:
        screen['type'] = 'completion'
    elif 'Key Takeaways' in section or 'Summary' in section.split('\n')[0]:
        screen['type'] = 'summary'
    elif 'Inclusive Practice' in section or 'Tips' in section.split('\n')[0]:
        screen['type'] = 'tips'

    # Extract narration
    narration_match = re.search(r'\*\*Narration.*?\*\*\s*\n\s*\|?\s*(.+?)(?:\s*\||\n\s*\|?\s*\*\*Audio)', section, re.DOTALL)
    if narration_match:
        screen['narration'] = narration_match.group(1).strip().replace('|', '').strip()

    # Extract setup text for scenarios
    setup_match = re.search(r'\*\*Setup Text:\*\*.*?\|\s*(.+?)(?:\s*\|\s*\n|\n\s*\|.*?Interaction)', section, re.DOTALL)
    if not setup_match:
        setup_match = re.search(r'Scenario — Branching\s*\|\s*(.+?)(?:\s*\|\s*(?:Realistic|Warm|Clean))', section, re.DOTALL)
    if setup_match:
        text = setup_match.group(1).strip()
        text = re.sub(r'\s*\|.*$', '', text, flags=re.MULTILINE).strip()
        screen['setup_text'] = text

    # Extract choices for scenarios
    choices = []
    choice_pattern = re.findall(
        r'\*\*Choice ([A-C])\s*\(([^)]+)\):\*\*\s*(.+?)\n\s*\|?\s*→\s*\*(.+?)\*',
        section, re.DOTALL
    )
    for letter, quality, text, consequence in choice_pattern:
        choices.append({
            'letter': letter,
            'quality': quality.strip(),
            'text': text.strip().replace('|', '').strip(),
            'consequence': consequence.strip().replace('|', '').strip(),
        })
    if choices:
        screen['choices'] = choices

    # Extract debrief
    debrief_match = re.search(r'\*\*Debrief:\*\*\s*(.+?)(?:\s*\||\n\s*\|?\s*\*\*(?:Audio|Principle))', section, re.DOTALL)
    if debrief_match:
        screen['debrief'] = debrief_match.group(1).strip().replace('|', '').strip()

    # Extract questions
    questions = []
    q_pattern = re.findall(
        r'\*\*Question \d+\s*\(([^)]+)\):\*\*\s*\n\s*\|?\s*(.+?)(?=\n\s*\|?\s*\*\*(?:Question|Audio|Narration)|\n---|\Z)',
        section, re.DOTALL
    )
    for q_type, q_body in q_pattern:
        q = {'type': q_type.strip(), 'body': q_body.strip()}

        # Parse question text and options
        lines = [l.strip().lstrip('|').strip() for l in q_body.strip().split('\n') if l.strip()]
        q['stem'] = lines[0] if lines else ''
        q['options'] = []
        for line in lines[1:]:
            opt_match = re.match(r'([a-d])\)\s*(.+?)(?:\s*✓)?$', line)
            if opt_match:
                is_correct = '✓' in line
                q['options'].append({
                    'letter': opt_match.group(1),
                    'text': opt_match.group(2).strip().rstrip('✓').strip(),
                    'correct': is_correct,
                })
        questions.append(q)
    if questions:
        screen['questions'] = questions

    # Extract feedback
    correct_match = re.search(r'Correct:\s*"?(.+?)"?\s*(?:\n|\|)', section)
    if correct_match:
        screen['feedback_correct'] = correct_match.group(1).strip().strip('"')
    incorrect_match = re.search(r'Incorrect:\s*"?(.+?)"?\s*(?:\n|\|)', section)
    if incorrect_match:
        screen['feedback_incorrect'] = incorrect_match.group(1).strip().strip('"')

    # Extract on-screen text items
    text_items = re.findall(r'•\s*(.+?)(?:\s*\||\s*$)', section, re.MULTILINE)
    screen['text_items'] = [t.strip() for t in text_items if t.strip() and not t.strip().startswith('Captions:')]

    # Extract image filename
    img_match = re.search(r'\*\*Image:\*\*\s*(.+?)(?:\s*\||\s*$)', section, re.MULTILINE)
    if img_match:
        screen['image'] = img_match.group(1).strip()

    # Extract alt text
    alt_match = re.search(r'Alt:\s*"(.+?)"', section)
    if alt_match:
        screen['alt_text'] = alt_match.group(1).strip()

    # Extract audio file
    audio_match = re.search(r'\*\*Audio:\*\*\s*(voiceover_[\d.]+\.mp3)', section)
    if audio_match:
        screen['audio'] = audio_match.group(1)

    # Extract principle
    principle_match = re.search(r'\*\*Principle:\*\*\s*(.+?)(?:\s*\||\s*$)', section, re.MULTILINE)
    if principle_match:
        screen['principle'] = principle_match.group(1).strip()

    return screen


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def create_title_slide(prs, info, guide_num):
    """Create the course title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])

    # Title
    title_ph = slide.placeholders[0]
    add_text_to_placeholder(title_ph, f"Guide {guide_num:02d}", font_size=28, bold=True, color=UHN_NAVY)

    # Subtitle
    subtitle_ph = slide.placeholders[11]
    title_text = info.get('title', 'Accessibility First').replace('Master Storyboard — Guide 01: ', '')
    lines = [
        title_text,
        "",
        "Accessibility First Course Series",
        "University Health Network | Toronto, Ontario",
        f"Seat Time: {info.get('seat_time', '15–20 min')}",
    ]
    add_multiline_text(subtitle_ph, lines, font_size=14, color=UHN_NAVY)

    return slide


def create_content_slide(prs, screen, layout_idx=LAYOUT_TEXT_BLUE):
    """Create a standard content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Title
    title_ph = slide.placeholders[0]
    screen_label = f"Screen {screen['number']}"
    add_text_to_placeholder(title_ph, f"{screen_label} — {screen['title']}", font_size=22, bold=True, color=UHN_NAVY)

    # Body content
    body_ph = slide.placeholders[15]
    body_lines = []

    # Add relevant text items
    for item in screen.get('text_items', []):
        if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:']):
            body_lines.append(f"• {item}")

    if body_lines:
        add_multiline_text(body_ph, body_lines, font_size=12, color=UHN_NAVY)

    # Add narration as notes
    if screen.get('narration'):
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = f"NARRATION:\n{screen['narration']}"
        if screen.get('audio'):
            notes_tf.text += f"\n\nAudio: {screen['audio']}"

    return slide


def create_scenario_slides(prs, screen):
    """Create scenario setup slide + consequence slides for each choice."""
    slides = []

    # --- Setup slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_BLUE])
    title_ph = slide.placeholders[0]
    add_text_to_placeholder(title_ph, f"Screen {screen['number']} — {screen['title']}", font_size=22, bold=True, color=UHN_NAVY)

    body_ph = slide.placeholders[15]
    tf = body_ph.text_frame
    tf.clear()
    tf.word_wrap = True

    # Setup text
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = screen.get('setup_text', 'Read the scenario below.')
    set_text_props(run, font_size=12, color=UHN_NAVY)

    # Add "What do you do?" prompt
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = "What do you do?"
    set_text_props(run2, font_size=14, bold=True, color=UHN_COBALT)

    # List choices
    for choice in screen.get('choices', []):
        p_choice = tf.add_paragraph()
        p_choice.space_before = Pt(8)
        run_label = p_choice.add_run()
        run_label.text = f"Choice {choice['letter']}: "
        set_text_props(run_label, font_size=12, bold=True, color=UHN_NAVY)
        run_text = p_choice.add_run()
        run_text.text = choice['text']
        set_text_props(run_text, font_size=12, color=UHN_NAVY)

    # Narration in notes
    if screen.get('narration'):
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"NARRATION:\n{screen['narration']}"

    slides.append(('setup', slide))

    # --- Consequence slide for each choice ---
    for choice in screen.get('choices', []):
        c_slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_NEUTRAL])

        c_title = c_slide.placeholders[0]
        quality = choice['quality']
        if 'Best' in quality:
            quality_color = CORRECT_GREEN
        elif 'Poor' in quality:
            quality_color = INCORRECT_RED
        else:
            quality_color = UHN_COBALT

        add_text_to_placeholder(
            c_title,
            f"Screen {screen['number']} — Choice {choice['letter']} ({quality})",
            font_size=20, bold=True, color=quality_color
        )

        c_body = c_slide.placeholders[15]
        tf = c_body.text_frame
        tf.clear()
        tf.word_wrap = True

        # What you chose
        p1 = tf.paragraphs[0]
        run_label = p1.add_run()
        run_label.text = "Your choice: "
        set_text_props(run_label, font_size=12, bold=True, color=UHN_NAVY)
        run_text = p1.add_run()
        run_text.text = choice['text']
        set_text_props(run_text, font_size=12, color=UHN_NAVY)

        # Consequence
        p2 = tf.add_paragraph()
        p2.space_before = Pt(12)
        run_label2 = p2.add_run()
        run_label2.text = "Result: "
        set_text_props(run_label2, font_size=12, bold=True, color=quality_color)
        run_cons = p2.add_run()
        run_cons.text = choice['consequence']
        set_text_props(run_cons, font_size=12, color=UHN_NAVY)

        # Debrief on the best choice slide
        if 'Best' in quality and screen.get('debrief'):
            p3 = tf.add_paragraph()
            p3.space_before = Pt(16)
            run_d_label = p3.add_run()
            run_d_label.text = "Key takeaway: "
            set_text_props(run_d_label, font_size=12, bold=True, color=UHN_TEAL)
            run_d = p3.add_run()
            run_d.text = screen['debrief']
            set_text_props(run_d, font_size=12, color=UHN_NAVY)

        # Principle
        if screen.get('principle') and 'Best' in quality:
            p4 = tf.add_paragraph()
            p4.space_before = Pt(12)
            run_p = p4.add_run()
            run_p.text = f"Principle: {screen['principle']}"
            set_text_props(run_p, font_size=10, italic=True, color=UHN_COBALT)

        slides.append((f'choice_{choice["letter"]}', c_slide))

    return slides


def create_question_slides(prs, screen):
    """Create question slide + correct feedback slide + incorrect feedback slide per question."""
    slides = []

    questions = screen.get('questions', [])
    if not questions:
        # Fallback: just create a content slide
        return [('content', create_content_slide(prs, screen))]

    for qi, q in enumerate(questions):
        q_num = qi + 1
        q_label = f"Screen {screen['number']}"
        if len(questions) > 1:
            q_label += f" — Q{q_num}"

        # --- Question slide ---
        q_slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_BLUE])
        q_title = q_slide.placeholders[0]
        add_text_to_placeholder(q_title, f"{q_label} — {screen['title']}", font_size=20, bold=True, color=UHN_NAVY)

        q_body = q_slide.placeholders[15]
        tf = q_body.text_frame
        tf.clear()
        tf.word_wrap = True

        # Question stem
        p_stem = tf.paragraphs[0]
        run_stem = p_stem.add_run()
        run_stem.text = q['stem']
        set_text_props(run_stem, font_size=14, bold=True, color=UHN_NAVY)

        # Options
        for opt in q.get('options', []):
            p_opt = tf.add_paragraph()
            p_opt.space_before = Pt(8)
            run_opt = p_opt.add_run()
            run_opt.text = f"    {opt['letter']})  {opt['text']}"
            set_text_props(run_opt, font_size=12, color=UHN_NAVY)

        slides.append(('question', q_slide))

        # --- Correct feedback slide ---
        correct_slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_NEUTRAL])
        c_title = correct_slide.placeholders[0]
        add_text_to_placeholder(c_title, f"{q_label} — Correct", font_size=20, bold=True, color=CORRECT_GREEN)

        c_body = correct_slide.placeholders[15]
        tf_c = c_body.text_frame
        tf_c.clear()
        tf_c.word_wrap = True

        # Show correct answer
        correct_opt = next((o for o in q.get('options', []) if o.get('correct')), None)
        p_ans = tf_c.paragraphs[0]
        run_check = p_ans.add_run()
        run_check.text = "✓ Correct!"
        set_text_props(run_check, font_size=18, bold=True, color=CORRECT_GREEN)

        if correct_opt:
            p_a = tf_c.add_paragraph()
            p_a.space_before = Pt(8)
            run_a = p_a.add_run()
            run_a.text = f"Answer: {correct_opt['letter']}) {correct_opt['text']}"
            set_text_props(run_a, font_size=12, bold=True, color=UHN_NAVY)

        feedback = screen.get('feedback_correct', 'Well done! You selected the correct answer.')
        p_fb = tf_c.add_paragraph()
        p_fb.space_before = Pt(12)
        run_fb = p_fb.add_run()
        run_fb.text = feedback
        set_text_props(run_fb, font_size=12, color=UHN_NAVY)

        slides.append(('correct', correct_slide))

        # --- Incorrect feedback slide ---
        incorrect_slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_NEUTRAL])
        i_title = incorrect_slide.placeholders[0]
        add_text_to_placeholder(i_title, f"{q_label} — Incorrect", font_size=20, bold=True, color=INCORRECT_RED)

        i_body = incorrect_slide.placeholders[15]
        tf_i = i_body.text_frame
        tf_i.clear()
        tf_i.word_wrap = True

        p_wrong = tf_i.paragraphs[0]
        run_x = p_wrong.add_run()
        run_x.text = "✗ Not quite."
        set_text_props(run_x, font_size=18, bold=True, color=INCORRECT_RED)

        feedback_inc = screen.get('feedback_incorrect', 'Review the content and try again.')
        p_fb_i = tf_i.add_paragraph()
        p_fb_i.space_before = Pt(12)
        run_fb_i = p_fb_i.add_run()
        run_fb_i.text = feedback_inc
        set_text_props(run_fb_i, font_size=12, color=UHN_NAVY)

        if correct_opt:
            p_correct_ans = tf_i.add_paragraph()
            p_correct_ans.space_before = Pt(12)
            run_ca = p_correct_ans.add_run()
            run_ca.text = f"The correct answer is: {correct_opt['letter']}) {correct_opt['text']}"
            set_text_props(run_ca, font_size=12, bold=True, color=CORRECT_GREEN)

        slides.append(('incorrect', incorrect_slide))

    return slides


def create_objectives_slide(prs, screen):
    """Create a learning objectives slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_BLUE])

    title_ph = slide.placeholders[0]
    add_text_to_placeholder(title_ph, f"Screen {screen['number']} — {screen['title']}", font_size=22, bold=True, color=UHN_NAVY)

    body_ph = slide.placeholders[15]
    # Extract numbered objectives from text items
    objectives = [item for item in screen.get('text_items', []) if any(item.startswith(f"{n}.") for n in range(1, 10))]
    if not objectives:
        objectives = [item for item in screen.get('text_items', []) if 'CLO' not in item and 'Audio' not in item]

    add_bullet_list(body_ph, objectives, font_size=13, color=UHN_NAVY)

    if screen.get('narration'):
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"NARRATION:\n{screen['narration']}"

    return slide


def create_reflection_slide(prs, screen):
    """Create a reflection/MAP slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_NEUTRAL])

    title_ph = slide.placeholders[0]
    add_text_to_placeholder(title_ph, f"Screen {screen['number']} — {screen['title']}", font_size=22, bold=True, color=UHN_TEAL)

    body_ph = slide.placeholders[15]
    tf = body_ph.text_frame
    tf.clear()
    tf.word_wrap = True

    # Add prompt or text items
    items = screen.get('text_items', [])
    reflection_items = [i for i in items if 'Audio' not in i and 'Captions' not in i and 'keyboard' not in i]

    for i, item in enumerate(reflection_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = item
        set_text_props(run, font_size=13, color=UHN_NAVY, italic=('private' in item.lower() or 'not submitted' in item.lower()))

    if screen.get('narration'):
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"NARRATION:\n{screen['narration']}"

    return slide


def create_summary_slide(prs, screen):
    """Create a key takeaways / summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TEXT_BLUE])

    title_ph = slide.placeholders[0]
    add_text_to_placeholder(title_ph, f"Screen {screen['number']} — {screen['title']}", font_size=22, bold=True, color=UHN_NAVY)

    body_ph = slide.placeholders[15]
    items = []
    for item in screen.get('text_items', []):
        clean = re.sub(r'^\d+\.\s*', '', item)
        if clean and 'Audio' not in clean and 'Captions' not in clean:
            items.append(clean)

    add_bullet_list(body_ph, items, font_size=12, color=UHN_NAVY)

    if screen.get('narration'):
        notes = slide.notes_slide.notes_text_frame
        notes.text = f"NARRATION:\n{screen['narration']}"

    return slide


def create_completion_slide(prs, screen, guide_num):
    """Create the final resources / completion slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CLOSING])
    # Closing slide layout has no placeholders — add text boxes manually
    from pptx.util import Inches

    # Completion message
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Guide {guide_num:02d} Complete"
    set_text_props(run, font_size=32, bold=True, color=UHN_WHITE)

    # Resources
    items = [i for i in screen.get('text_items', []) if 'Audio' not in i and 'Captions' not in i and 'badge' not in i.lower()]
    if items:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, item in enumerate(items[:6]):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            set_text_props(run, font_size=11, color=UHN_WHITE)

    return slide


# ---------------------------------------------------------------------------
# Main generation
# ---------------------------------------------------------------------------

def generate_ppt(guide_num: int, dry_run: bool = False):
    """Generate the full PPT from the master storyboard."""
    guide_folder = BUILD_OUTPUT / GUIDE_FOLDERS[guide_num]
    storyboard_path = guide_folder / "01-master-storyboard" / f"MASTER-STORYBOARD-GUIDE-{guide_num:02d}.md"

    if not storyboard_path.exists():
        print(f"Error: Storyboard not found: {storyboard_path}")
        sys.exit(1)

    if not TEMPLATE_PATH.exists():
        print(f"Error: UHN PPT template not found: {TEMPLATE_PATH}")
        sys.exit(1)

    print(f"Parsing storyboard: {storyboard_path.name}")
    data = parse_storyboard(storyboard_path)
    info = data['info']
    screens = data['screens']

    print(f"Found {len(screens)} screens")
    print(f"Course: {info.get('title', 'Unknown')}")
    print(f"Seat time: {info.get('seat_time', 'N/A')}")
    print()

    # Plan slides
    slide_plan = []
    slide_plan.append(('title', 'Course Title Slide', None))

    for screen in screens:
        stype = screen['type']
        label = f"Screen {screen['number']} — {screen['title']}"

        if stype == 'scenario':
            slide_plan.append(('scenario_setup', f"{label} (Setup)", screen))
            for choice in screen.get('choices', []):
                slide_plan.append(('scenario_choice', f"  → Choice {choice['letter']} ({choice['quality']})", screen))
        elif stype == 'question':
            questions = screen.get('questions', [])
            for qi, q in enumerate(questions):
                q_label = f"{label}"
                if len(questions) > 1:
                    q_label += f" Q{qi+1}"
                slide_plan.append(('question', f"{q_label} (Question)", screen))
                slide_plan.append(('correct', f"  → Correct Feedback", screen))
                slide_plan.append(('incorrect', f"  → Incorrect Feedback", screen))
        else:
            slide_plan.append((stype, label, screen))

    print(f"Slide plan ({len(slide_plan)} slides):")
    print("-" * 60)
    for i, (stype, label, _) in enumerate(slide_plan, 1):
        print(f"  {i:3d}. [{stype:16s}] {label}")
    print("-" * 60)
    print()

    if dry_run:
        print("Dry run — no PPT generated.")
        return

    # Generate PPT
    print("Generating PPT from UHN template...")
    prs = Presentation(str(TEMPLATE_PATH))

    # Remove any existing slides from the template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 1. Title slide
    create_title_slide(prs, info, guide_num)
    print("  ✓ Title slide")

    # 2. Screen slides
    for screen in screens:
        stype = screen['type']
        label = f"Screen {screen['number']} — {screen['title']}"

        if stype == 'welcome':
            create_content_slide(prs, screen)
            print(f"  ✓ {label} (welcome)")

        elif stype == 'objectives':
            create_objectives_slide(prs, screen)
            print(f"  ✓ {label} (objectives)")

        elif stype == 'scenario':
            scenario_slides = create_scenario_slides(prs, screen)
            print(f"  ✓ {label} (scenario: {len(scenario_slides)} slides)")

        elif stype == 'question':
            q_slides = create_question_slides(prs, screen)
            print(f"  ✓ {label} (questions: {len(q_slides)} slides)")

        elif stype == 'reflection' or stype == 'action_planning':
            create_reflection_slide(prs, screen)
            print(f"  ✓ {label} ({stype})")

        elif stype == 'summary' or stype == 'tips':
            create_summary_slide(prs, screen)
            print(f"  ✓ {label} ({stype})")

        elif stype == 'completion':
            create_completion_slide(prs, screen, guide_num)
            print(f"  ✓ {label} (completion)")

        else:
            create_content_slide(prs, screen)
            print(f"  ✓ {label} (content)")

    # Save
    output_dir = guide_folder / "02-ppt-storyboard"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"UHN_Guide{guide_num:02d}_Storyboard.pptx"
    prs.save(str(output_path))

    total_slides = len(prs.slides)
    print(f"\n✓ PPT saved: {output_path}")
    print(f"  Total slides: {total_slides}")
    print(f"  Template: {TEMPLATE_PATH.name}")
    print(f"  Size: {output_path.stat().st_size / 1024:.0f} KB")


def main():
    parser = argparse.ArgumentParser(description="Generate PPT storyboard for UHN Accessibility First courses")
    parser.add_argument("--guide", type=int, required=True, help="Guide number (1-18)")
    parser.add_argument("--dry-run", action="store_true", help="Preview slide plan without generating")
    args = parser.parse_args()

    if args.guide not in GUIDE_FOLDERS:
        print(f"Error: Guide {args.guide} not found. Valid: 1-18")
        sys.exit(1)

    generate_ppt(args.guide, dry_run=args.dry_run)


if __name__ == "__main__":
    main()
