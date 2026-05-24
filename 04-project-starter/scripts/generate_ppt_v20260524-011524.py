#!/usr/bin/env python3
"""
PPT Storyboard Generator for UHN Accessibility First Course Series
Converts Master Storyboard markdown into branded PowerPoint for Storyline 360 import.

Builds every slide from scratch using python-pptx shapes on a blank layout,
matching the HTML mockup design system (1920x1080 / 16:9).

Each screen becomes one or more slides:
- Content screens → 1 slide
- Question screens → question slide + correct feedback slide + incorrect feedback slide
- Branching scenarios → setup slide + consequence slides per choice (A/B/C) + debrief slide

Usage:
    python generate_ppt.py --guide 1
    python generate_ppt.py --guide 1 --dry-run    # Preview slide plan without generating

Requirements:
    pip install python-pptx Pillow
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("Error: 'python-pptx' package not installed. Run: pip install python-pptx")
    sys.exit(1)

BASE_DIR = Path(__file__).resolve().parent.parent.parent
BUILD_OUTPUT = BASE_DIR / "05-build-output"
TEMPLATE_PATH = BASE_DIR / "02-branding-and-style" / "uhn-templates" / "UHN_Presentation_Template_Primary (1).pptx"
LOGO_PATH = BASE_DIR / "02-branding-and-style" / "logos" / "uhn_logo.png"

GUIDE_FOLDERS = {
    1: "01-Foundations-of-Disability-Inclusion-and-Accessible-Design",
    2: "02-Perceptions-Attitudes-and-Barriers",
    3: "03-Vision-Disabilities",
    4: "04-Sensory-Hearing-and-Communication-Disabilities",
    5: "05-Physical-Disabilities-and-Mobility",
    6: "06-Mental-Health-Disabilities",
    7: "07-Intellectual-Developmental-and-Learning-Disabilities",
    8: "08-Non-Visible-Disabilities",
    9: "09-Aging-Disability-and-Intersectionality",
    10: "10-Engaging-with-Confidence-and-Respect",
    11: "11-Service-Animals-Guide-Dogs-and-Non-Service-Animals",
    12: "12-Support-Persons",
    13: "13-Assistive-Devices",
    14: "14-Communication-and-Information-Accessibility",
    15: "15-Neurodiversity-and-Sensory-Regulation",
    16: "16-Trauma-Informed-Accessibility",
    17: "17-Accessibility-in-Crisis-Situations-and-De-escalation",
    18: "18-Indigenous-Peoples-and-Accessibility",
}

# Guide subtitle lookup
GUIDE_SUBTITLES = {
    1: "Foundations",
    2: "Perceptions & Attitudes",
    3: "Vision Disabilities",
    4: "Hearing & Communication",
    5: "Physical & Mobility",
    6: "Mental Health",
    7: "Intellectual & Developmental",
    8: "Non-Visible Disabilities",
    9: "Aging & Intersectionality",
    10: "Confidence & Respect",
    11: "Service Animals",
    12: "Support Persons",
    13: "Assistive Devices",
    14: "Communication Access",
    15: "Neurodiversity",
    16: "Trauma-Informed",
    17: "Crisis & De-escalation",
    18: "Indigenous Peoples",
}

# ---------------------------------------------------------------------------
# Design system colours (from HTML mockup CSS variables)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x19, 0x28, 0x58)
NAVY_DEEP = RGBColor(0x0F, 0x1A, 0x3D)
RED = RGBColor(0xC0, 0x23, 0x3B)
RED_BG = RGBColor(0xFD, 0xF2, 0xF3)
COBALT = RGBColor(0x24, 0x5B, 0xAA)
COBALT_BG = RGBColor(0xEA, 0xF1, 0xFA)
LILAC = RGBColor(0xC4, 0x8A, 0xBD)
LILAC_BG = RGBColor(0xF7, 0xED, 0xF5)
CHARTREUSE = RGBColor(0x74, 0xAE, 0x54)
CHARTREUSE_BG = RGBColor(0xEE, 0xF6, 0xE7)
NEUTRAL = RGBColor(0xE4, 0xE4, 0xE4)
NEUTRAL_2 = RGBColor(0xF4, 0xF4, 0xF4)
INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Common chrome dimensions
TOPBAR_H = Inches(0.583)        # 84px
CONTENT_PAD_X = Inches(0.667)   # 96px
CONTENT_TOP = Inches(0.75)      # below topbar
FOOTER_Y = Inches(6.95)
FOOTER_H = Inches(0.417)        # ~60px


# ---------------------------------------------------------------------------
# Helper: text formatting
# ---------------------------------------------------------------------------

def _set_run(run, text, font_name="Arial", size=14, bold=False, italic=False, color=None):
    """Apply text formatting to a run."""
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def _add_paragraph(tf, text, font_name="Arial", size=14, bold=False, italic=False,
                   color=None, alignment=None, space_before=0, space_after=0):
    """Add a paragraph with a single run to a text frame."""
    p = tf.add_paragraph()
    if alignment:
        p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    run = p.add_run()
    _set_run(run, text, font_name=font_name, size=size, bold=bold, italic=italic, color=color)
    return p


def _init_textbox(slide, left, top, width, height, word_wrap=True):
    """Add a textbox and return the text frame (first paragraph cleared)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    # Clear the default empty paragraph
    tf.paragraphs[0].clear()
    return tf, txBox


def _add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # no border by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()  # transparent
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = Pt(border_width)
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
    return shape


# ---------------------------------------------------------------------------
# Common chrome: topbar + footer
# ---------------------------------------------------------------------------

def _add_topbar(slide, guide_num, section_name=""):
    """Add the navy topbar with UHN logo + breadcrumbs."""
    # Navy background bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, TOPBAR_H, fill_color=NAVY)

    # UHN logo
    logo_h = Inches(0.333)  # 48px
    logo_top = Inches((0.583 - 0.333) / 2)  # centered in topbar
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.333), logo_top, height=logo_h)

    # Separator line (thin vertical)
    if LOGO_PATH.exists():
        sep_x = Inches(1.5)
    else:
        sep_x = Inches(0.333)
    _add_rect(slide, sep_x, Inches(0.125), Inches(0.014), Inches(0.333), fill_color=RGBColor(0x4A, 0x5A, 0x8A))

    # "Accessibility First" label
    tf, _ = _init_textbox(slide, sep_x + Inches(0.15), Inches(0.05), Inches(2.5), TOPBAR_H)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    _set_run(run, "Accessibility First", font_name="Arial Black", size=11, bold=True, color=WHITE)

    # Breadcrumbs on the right
    subtitle = GUIDE_SUBTITLES.get(guide_num, "")
    breadcrumb = f"Guide {guide_num:02d} \u00b7 {subtitle}"
    if section_name:
        breadcrumb += f"  >  {section_name}"
    bc_tf, _ = _init_textbox(slide, Inches(8.5), Inches(0.05), Inches(4.5), TOPBAR_H)
    bc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = bc_tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    _set_run(run, breadcrumb, font_name="Arial", size=9, color=RGBColor(0xAA, 0xBB, 0xDD))


def _add_footer(slide, slide_index, total_slides):
    """Add footer with ACCESSIBILITY FIRST tag + progress dots."""
    # Red dot
    _add_rect(slide, Inches(0.667), FOOTER_Y + Inches(0.08), Inches(0.083), Inches(0.083), fill_color=RED)

    # "ACCESSIBILITY FIRST" tag
    tf, _ = _init_textbox(slide, Inches(0.83), FOOTER_Y, Inches(2.5), Inches(0.35))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    _set_run(run, "ACCESSIBILITY FIRST", font_name="Arial Black", size=7, bold=True, color=NAVY)

    # Progress dots on the right
    dot_w = Inches(0.18)
    dot_h = Inches(0.055)
    dot_gap = Inches(0.07)
    # Max 20 dots to avoid overcrowding
    max_dots = min(total_slides, 20)
    step = max(1, total_slides // max_dots)
    dots_to_show = (total_slides + step - 1) // step
    start_x = Inches(13.333) - CONTENT_PAD_X - (dots_to_show * (dot_w + dot_gap))
    dot_y = FOOTER_Y + Inches(0.07)

    for i in range(dots_to_show):
        actual_idx = i * step
        x = start_x + i * (dot_w + dot_gap)
        if actual_idx < slide_index:
            color = NAVY
        elif actual_idx == slide_index:
            color = RED
        else:
            color = NEUTRAL
        _add_rect(slide, x, dot_y, dot_w, dot_h, fill_color=color)


def _add_eyebrow(slide, text, y=None):
    """Add an eyebrow label (red, uppercase, letter-spaced)."""
    if y is None:
        y = CONTENT_TOP
    tf, _ = _init_textbox(slide, CONTENT_PAD_X, y, Inches(10), Inches(0.35))
    run = tf.paragraphs[0].add_run()
    # Simulate letter-spacing with spaces
    spaced = "  ".join(text.upper())
    _set_run(run, spaced, font_name="Arial Black", size=9, bold=True, color=RED)
    return y + Inches(0.35)


def _add_slide_title(slide, text, y=None):
    """Add the main slide H1 title."""
    if y is None:
        y = CONTENT_TOP + Inches(0.35)
    tf, _ = _init_textbox(slide, CONTENT_PAD_X, y, Inches(11), Inches(0.7))
    run = tf.paragraphs[0].add_run()
    _set_run(run, text, font_name="Arial Black", size=36, bold=True, color=NAVY)
    return y + Inches(0.7)


def _add_lede(slide, text, y=None):
    """Add a lede / subtitle paragraph."""
    if y is None:
        y = CONTENT_TOP + Inches(1.1)
    tf, _ = _init_textbox(slide, CONTENT_PAD_X, y, Inches(10), Inches(0.6))
    run = tf.paragraphs[0].add_run()
    _set_run(run, text, font_name="Arial", size=15, color=MUTED)
    return y + Inches(0.6)


def _add_notes(slide, screen):
    """Add narration + audio info to slide notes."""
    notes_text = ""
    if screen.get('narration'):
        notes_text = f"NARRATION:\n{screen['narration']}"
    if screen.get('audio'):
        notes_text += f"\n\nAudio: {screen['audio']}"
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def _get_blank_layout(prs):
    """Get the blank slide layout (or the closest one with fewest placeholders)."""
    # Try to find a layout named "Blank"
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            return layout
    # Fallback: use the layout with fewest placeholders
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _add_blank_slide(prs):
    """Add a blank slide with white background."""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    # Clear any placeholder shapes inherited from layout
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    # Set white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


# ---------------------------------------------------------------------------
# Storyboard parser (UNCHANGED)
# ---------------------------------------------------------------------------

def parse_storyboard(filepath: Path) -> dict:
    """Parse the master storyboard markdown into structured data."""
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    # Extract header info
    info = {}
    title_match = re.search(r'^# (.+)', content, re.MULTILINE)
    if title_match:
        info['title'] = title_match.group(1).strip()

    seat_match = re.search(r'Seat Time\s*\|\s*(.+?)(?:\s*\|)', content)
    if seat_match:
        info['seat_time'] = seat_match.group(1).strip()

    screens_match = re.search(r'Screens\s*\|\s*(\d+)', content)
    if screens_match:
        info['total_screens'] = int(screens_match.group(1))

    clo_match = re.search(r'\*\*CLOs?:\*\*\s*(.+?)$', content, re.MULTILINE)
    if clo_match:
        info['clos'] = clo_match.group(1).strip()

    # Parse individual screens
    screen_sections = re.split(r'(?=^## Screen \d)', content, flags=re.MULTILINE)
    screens = []

    for section in screen_sections:
        if not section.strip().startswith('## Screen'):
            continue

        screen = parse_screen(section)
        if screen:
            screens.append(screen)

    return {'info': info, 'screens': screens}


def flatten_table_column(section: str, col_idx: int) -> str:
    """Extract all content from a specific column of a markdown table.
    col_idx: 0=Step, 1=Activities, 2=Design Guide"""
    lines = []
    for line in section.split('\n'):
        line = line.strip()
        if not line.startswith('|') or line.startswith('|---') or line.startswith('| Step'):
            continue
        cells = [c.strip() for c in line.split('|')]
        # cells[0] is empty (before first |), cells[-1] is empty (after last |)
        if len(cells) > col_idx + 1:
            cell = cells[col_idx + 1].strip()
            if cell:
                lines.append(cell)
    return '\n'.join(lines)


def parse_screen(section: str) -> dict:
    """Parse a single screen section into structured data."""
    screen = {}

    # Screen number and title
    header_match = re.search(r'## Screen ([\d.]+)\s*[—–-]\s*(.+)', section)
    if header_match:
        screen['number'] = header_match.group(1).strip()
        screen['title'] = header_match.group(2).strip()

    # Flatten the Activities column (col 1) for easier parsing
    activities = flatten_table_column(section, 1)
    design_col = flatten_table_column(section, 2)

    # Determine screen type based on title and content
    title = screen.get('title', '')
    screen['type'] = 'content'  # default

    if 'Welcome' in title:
        screen['type'] = 'welcome'
    elif 'Learning Objectives' in title or 'Objectives' in title:
        screen['type'] = 'objectives'
    elif 'Decision Tree' in title and ('Choice A' in activities or 'Choice B' in activities):
        screen['type'] = 'decision_tree'
    elif 'Scenario' in title and ('Branching' in activities or 'Choice A' in activities):
        screen['type'] = 'scenario'
    elif 'Knowledge Check' in title or 'Assessment' in title or 'Question' in activities[:50]:
        screen['type'] = 'question'
    elif 'Listen and Reflect' in title or 'Podcast' in title:
        screen['type'] = 'podcast'
    elif 'Series Progress' in title or 'Progress Map' in title:
        screen['type'] = 'progress'
    elif 'Reflection' in title:
        screen['type'] = 'reflection'
    elif 'MAP' in title or 'Action Planning' in title:
        screen['type'] = 'action_planning'
    elif 'Completion' in title or 'Resources' in title:
        screen['type'] = 'completion'
    elif 'Key Takeaways' in title or 'Summary' in title:
        screen['type'] = 'summary'
    elif 'Impact' in title:
        screen['type'] = 'impact'
    elif 'Inclusive Practice' in title or 'Tips' in title:
        screen['type'] = 'tips'

    # Extract narration — find the narration text in activities column
    narration_match = re.search(
        r'\*\*Narration\s*\([^)]*\):\*\*\s*\n?(.*?)(?:\*\*Audio|\*\*Refs|\*\*Audio:|\Z)',
        activities, re.DOTALL
    )
    if narration_match:
        narration = narration_match.group(1).strip()
        # Clean up any remaining markdown
        narration = re.sub(r'\*\*[^*]+\*\*', '', narration).strip()
        if narration:
            screen['narration'] = narration

    # Extract setup text for scenarios (text before "What do you do?")
    if screen['type'] == 'scenario':
        setup_match = re.search(
            r'(?:Setup Text:|Scenario — Branching)\s*\n?(.*?)(?:\*\*What do you do\?\*\*|What do you do\?)',
            activities, re.DOTALL
        )
        if not setup_match:
            # Try getting the long text block from activities
            lines = activities.split('\n')
            setup_lines = []
            for line in lines:
                if 'Choice A' in line or 'What do you do' in line:
                    break
                clean = line.strip().lstrip('*').rstrip('*').strip()
                if clean and not any(x in clean for x in ['Setup Text:', 'Scenario', 'Interaction:', 'Time:', 'SME:', 'Screen 1.']):
                    setup_lines.append(clean)
            if setup_lines:
                screen['setup_text'] = ' '.join(setup_lines)
        else:
            text = setup_match.group(1).strip()
            # Remove markdown formatting artifacts
            text = re.sub(r'\*\*[^*]+:\*\*', '', text).strip()
            text = re.sub(r'\n+', ' ', text).strip()
            screen['setup_text'] = text

    # Extract choices for scenarios
    choices = []
    choice_pattern = re.finditer(
        r'\*\*Choice ([A-C])\s*\(([^)]+)\):\*\*\s*(.*?)→\s*\*(.*?)\*',
        activities, re.DOTALL
    )
    for m in choice_pattern:
        letter, quality, text, consequence = m.group(1), m.group(2), m.group(3), m.group(4)
        choices.append({
            'letter': letter.strip(),
            'quality': quality.strip(),
            'text': re.sub(r'\n+', ' ', text).strip(),
            'consequence': re.sub(r'\n+', ' ', consequence).strip(),
        })
    if choices:
        screen['choices'] = choices

    # Extract debrief
    debrief_match = re.search(r'\*\*Debrief:\*\*\s*(.*?)(?:\*\*Audio|\*\*Principle|\Z)', activities, re.DOTALL)
    if debrief_match:
        screen['debrief'] = re.sub(r'\n+', ' ', debrief_match.group(1)).strip()

    # Extract questions (for knowledge check screens)
    questions = []
    q_iter = re.finditer(
        r'\*\*Question \d+\s*\(([^)]+)\):\*\*\s*\n?(.*?)(?=\*\*Question \d+|\*\*Audio|\Z)',
        activities, re.DOTALL
    )
    for m in q_iter:
        q_type = m.group(1).strip()
        q_body = m.group(2).strip()
        q = {'type': q_type}

        # Parse question text and options from the body
        lines = [l.strip() for l in q_body.split('\n') if l.strip()]
        # First line is the stem
        stem_lines = []
        options = []
        for line in lines:
            opt_match = re.match(r'^([a-d])\)\s*(.+?)(\s*\u2713)?$', line)
            if opt_match:
                options.append({
                    'letter': opt_match.group(1),
                    'text': opt_match.group(2).strip(),
                    'correct': bool(opt_match.group(3)),
                })
            elif not options:  # still building the stem
                stem_lines.append(line)

        q['stem'] = ' '.join(stem_lines)
        q['options'] = options
        questions.append(q)

    if questions:
        screen['questions'] = questions

    # Extract feedback from design column
    correct_match = re.search(r'Correct:\s*"([^"]+)"', design_col)
    if correct_match:
        screen['feedback_correct'] = correct_match.group(1).strip()
    incorrect_match = re.search(r'Incorrect:\s*"([^"]+)"', design_col)
    if incorrect_match:
        screen['feedback_incorrect'] = incorrect_match.group(1).strip()

    # Extract on-screen text items from activities column (• bullets)
    text_items = re.findall(r'\u2022\s*(.+?)$', activities, re.MULTILINE)
    screen['text_items'] = [t.strip() for t in text_items if t.strip()]

    # Extract numbered list items (1. xxx, 2. xxx)
    numbered_items = re.findall(r'^(\d+)\.\s+(.+?)$', activities, re.MULTILINE)
    screen['numbered_items'] = [(num, text.strip()) for num, text in numbered_items if text.strip()]

    # Extract bold-label items (**Label:** text) — for models, frameworks, tips, steps
    skip_labels = {'Narration', 'Audio', 'Setup Text', 'On-Screen Text', 'Debrief',
                   'Refs', 'Prompt', 'Completion message', 'Next in series',
                   'Resources', 'Screen', 'What do you do?'}
    labeled_items = re.findall(r'\*\*([^*]+?):\*\*\s*(.+?)$', activities, re.MULTILINE)
    screen['labeled_items'] = [
        (label.strip(), text.strip()) for label, text in labeled_items
        if label.strip() not in skip_labels
        and not label.startswith('Question')
        and not label.startswith('Choice')
        and not label.startswith('Narration')
        and not label.startswith('Tips (')
        and not label.startswith('Summary (')
        and not label.startswith('My Action Planning')
        and not re.match(r'Screen \d', label)
    ]

    # Re-extract Tip items that may have been nested inside a parent match
    tip_items = re.findall(r'\*\*Tip (\d+):\*\*\s*(.+?)$', activities, re.MULTILINE)
    if tip_items and not any(l.startswith('Tip 1') for l, _ in screen['labeled_items']):
        screen['labeled_items'] = [
            (f'Tip {num}', text.strip()) for num, text in tip_items
        ]

    # Extract bold step items (**1. Label** — text) — for decision path
    step_items = re.findall(r'\*\*(\d+)\.\s*([^*]+)\*\*\s*[—–-]\s*(.+?)$', activities, re.MULTILINE)
    screen['step_items'] = [(num, label.strip(), text.strip()) for num, label, text in step_items]

    # Extract image filename from design column
    img_match = re.search(r'\*\*Image:\*\*\s*(.+?)$', design_col, re.MULTILINE)
    if img_match:
        screen['image'] = img_match.group(1).strip()

    # Extract alt text from design column
    alt_match = re.search(r'Alt:\s*"(.+?)"', design_col)
    if alt_match:
        screen['alt_text'] = alt_match.group(1).strip()

    # Extract audio file
    audio_match = re.search(r'\*\*Audio:\*\*\s*(voiceover_[\d.]+\.mp3)', activities)
    if audio_match:
        screen['audio'] = audio_match.group(1)

    # Extract principle from design column
    principle_match = re.search(r'\*\*Principle:\*\*\s*(.+?)$', design_col, re.MULTILINE)
    if principle_match:
        screen['principle'] = principle_match.group(1).strip()

    return screen


# ---------------------------------------------------------------------------
# Slide builders — mockup-matched, all from scratch
# ---------------------------------------------------------------------------

def _extract_section_name(screen):
    """Try to derive section name from screen number/title."""
    title = screen.get('title', '')
    # Strip markdown artifacts
    title = re.sub(r'\*+', '', title).strip()
    return title


def create_title_slide(prs, info, guide_num, total_slides):
    """Slide 1: Title slide (s-title) — Navy background, 2-column."""
    slide = _add_blank_slide(prs)

    # Navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Lilac accent line at top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), fill_color=LILAC)

    # --- LEFT COLUMN (x=0.667 to ~6.5) ---
    left_x = CONTENT_PAD_X

    # UHN logo (top-left)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left_x, Inches(0.5), height=Inches(0.5))

    # Series tag: "Guide XX of 18 · Foundations"
    subtitle = GUIDE_SUBTITLES.get(guide_num, "")
    series_text = f"Guide {guide_num:02d} of 18  \u00b7  {subtitle}"
    tf, _ = _init_textbox(slide, left_x, Inches(1.3), Inches(5.5), Inches(0.4))
    run = tf.paragraphs[0].add_run()
    _set_run(run, series_text, font_name="Arial", size=14, color=LILAC)

    # Main title (H1)
    raw_title = info.get('title', 'Accessibility First')
    # Strip "Master Storyboard — Guide XX: " prefix if present
    clean_title = re.sub(r'Master Storyboard\s*[—–-]\s*Guide\s*\d+:\s*', '', raw_title).strip()
    tf, _ = _init_textbox(slide, left_x, Inches(1.9), Inches(5.5), Inches(1.8))
    run = tf.paragraphs[0].add_run()
    _set_run(run, clean_title, font_name="Arial Black", size=42, bold=True, color=WHITE)

    # Subtitle
    tf, _ = _init_textbox(slide, left_x, Inches(3.8), Inches(5.5), Inches(0.6))
    run = tf.paragraphs[0].add_run()
    _set_run(run, "Accessibility First Course Series", font_name="Arial", size=16, color=RGBColor(0xCC, 0xCC, 0xDD))

    # Meta row
    seat_time = info.get('seat_time', '15-20 min')
    total_screens = info.get('total_screens', '')
    meta_items = [
        f"Duration: {seat_time}",
        "Format: eLearning (Storyline 360)",
        "Audience: All UHN staff",
    ]
    meta_y = Inches(4.7)
    for item in meta_items:
        tf, _ = _init_textbox(slide, left_x, meta_y, Inches(5), Inches(0.3))
        run = tf.paragraphs[0].add_run()
        _set_run(run, item, font_name="Arial", size=10, color=MUTED)
        meta_y += Inches(0.3)

    # --- RIGHT COLUMN (x=7.0 to ~12.5) ---
    right_x = Inches(7.0)
    right_w = Inches(5.667)

    # Photo placeholder
    photo_rect = _add_rounded_rect(slide, right_x, Inches(0.8), right_w, Inches(4.5),
                                    fill_color=RGBColor(0x2A, 0x3A, 0x6A))
    tf2, _ = _init_textbox(slide, right_x + Inches(1.5), Inches(2.7), Inches(3), Inches(0.5))
    run2 = tf2.paragraphs[0].add_run()
    _set_run(run2, "[Hero Photo Placeholder]", font_name="Arial", size=12, color=MUTED)

    # BEGIN button
    btn = _add_rounded_rect(slide, right_x + Inches(1.5), Inches(5.8), Inches(2.667), Inches(0.55),
                             fill_color=RED)
    tf_btn, _ = _init_textbox(slide, right_x + Inches(1.5), Inches(5.8), Inches(2.667), Inches(0.55))
    tf_btn.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_btn.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_btn = p.add_run()
    _set_run(run_btn, "BEGIN", font_name="Arial Black", size=14, bold=True, color=WHITE)

    # Footer (no progress dots on title)
    _add_rect(slide, Inches(0), FOOTER_Y, SLIDE_W, Inches(0.05), fill_color=LILAC)

    return slide


def create_objectives_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 2: Learning Objectives (s-obj) — 2x2 grid of objective cards."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    # Eyebrow
    section_num = screen.get('number', '2').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 What you will learn")

    # Title
    y = _add_slide_title(slide, "Learning Objectives", y)

    # Objective cards (2x2 grid) — prefer numbered_items
    objectives = [text for _, text in screen.get('numbered_items', [])]
    if not objectives:
        objectives = [item for item in screen.get('text_items', [])
                      if any(item.startswith(f"{n}.") for n in range(1, 10))]
    if not objectives:
        objectives = [item for item in screen.get('text_items', [])
                      if 'CLO' not in item and 'Audio' not in item]

    card_w = Inches(5.5)
    card_h = Inches(1.3)
    start_x = CONTENT_PAD_X
    start_y = Inches(2.8)
    gap_x = Inches(0.5)
    gap_y = Inches(0.35)

    for i, obj_text in enumerate(objectives[:4]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y_card = start_y + row * (card_h + gap_y)

        # Card background
        _add_rect(slide, x, y_card, card_w, card_h, fill_color=WHITE, border_color=NEUTRAL, border_width=0.5)
        # Navy top border accent
        _add_rect(slide, x, y_card, card_w, Inches(0.04), fill_color=NAVY)

        # Number
        tf_num, _ = _init_textbox(slide, x + Inches(0.15), y_card + Inches(0.12), Inches(0.5), Inches(0.5))
        run_num = tf_num.paragraphs[0].add_run()
        _set_run(run_num, str(i + 1), font_name="Arial Black", size=28, bold=True, color=NAVY)

        # Objective text — split title and description on " — "
        clean = re.sub(r'^\d+\.\s*', '', obj_text).strip()
        if ' — ' in clean:
            obj_title, obj_desc = clean.split(' — ', 1)
        elif ' - ' in clean:
            obj_title, obj_desc = clean.split(' - ', 1)
        else:
            obj_title, obj_desc = clean, ''

        # Title (bold)
        tf_title, _ = _init_textbox(slide, x + Inches(0.7), y_card + Inches(0.12), card_w - Inches(0.85), Inches(0.35))
        run_t = tf_title.paragraphs[0].add_run()
        _set_run(run_t, obj_title, font_name="Arial Black", size=12, bold=True, color=NAVY)

        # Description
        if obj_desc:
            tf_desc, _ = _init_textbox(slide, x + Inches(0.7), y_card + Inches(0.5), card_w - Inches(0.85), card_h - Inches(0.6))
            run_d = tf_desc.paragraphs[0].add_run()
            _set_run(run_d, obj_desc, font_name="Arial", size=10, color=MUTED)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides):
    """Generic content slide with topbar, eyebrow, title, and content cards/bullets."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 {screen.get('title', '')}")
    y = _add_slide_title(slide, screen.get('title', 'Content'), y)

    labeled = screen.get('labeled_items', [])
    step_items = screen.get('step_items', [])
    numbered = screen.get('numbered_items', [])
    body_items = [item for item in screen.get('text_items', [])
                  if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:'])]
    has_image = bool(screen.get('image'))

    # Determine content width based on image
    content_w = Inches(6.0) if has_image else Inches(11.5)
    content_x = CONTENT_PAD_X

    # --- Labeled items as cards (for models, quadrants, etc.) ---
    if labeled and len(labeled) >= 3:
        content_y = y + Inches(0.4)
        n_items = len(labeled)

        if n_items == 4:
            # 2x2 grid of cards (for AiP quadrants)
            card_w = content_w / 2 - Inches(0.15)
            card_h = Inches(1.6)
            for i, (label, text) in enumerate(labeled[:4]):
                col = i % 2
                row = i // 2
                cx = content_x + col * (card_w + Inches(0.3))
                cy = content_y + row * (card_h + Inches(0.25))

                _add_rect(slide, cx, cy, card_w, card_h, fill_color=WHITE, border_color=NEUTRAL, border_width=0.5)
                _add_rect(slide, cx, cy, card_w, Inches(0.04), fill_color=COBALT)

                # Label
                tf_l, _ = _init_textbox(slide, cx + Inches(0.2), cy + Inches(0.15), card_w - Inches(0.4), Inches(0.35))
                run_l = tf_l.paragraphs[0].add_run()
                _set_run(run_l, label, font_name="Arial Black", size=14, bold=True, color=NAVY)

                # Text
                tf_t, _ = _init_textbox(slide, cx + Inches(0.2), cy + Inches(0.55), card_w - Inches(0.4), card_h - Inches(0.7))
                run_t = tf_t.paragraphs[0].add_run()
                _set_run(run_t, text, font_name="Arial", size=10, color=INK)
        elif n_items == 3:
            # 3-column cards (for disability models)
            card_w = content_w / 3 - Inches(0.2)
            card_h = Inches(2.5)
            colors = [RED, COBALT, NAVY]
            for i, (label, text) in enumerate(labeled[:3]):
                cx = content_x + i * (card_w + Inches(0.3))
                cy = content_y

                _add_rect(slide, cx, cy, card_w, card_h, fill_color=WHITE, border_color=NEUTRAL, border_width=0.5)
                # Colored top bar
                _add_rect(slide, cx, cy, card_w, Inches(0.55), fill_color=colors[i % 3])

                # Label on colored bar
                tf_l, _ = _init_textbox(slide, cx + Inches(0.15), cy + Inches(0.08), card_w - Inches(0.3), Inches(0.4))
                run_l = tf_l.paragraphs[0].add_run()
                _set_run(run_l, label, font_name="Arial Black", size=13, bold=True, color=WHITE)

                # Text below
                tf_t, _ = _init_textbox(slide, cx + Inches(0.15), cy + Inches(0.65), card_w - Inches(0.3), card_h - Inches(0.8))
                run_t = tf_t.paragraphs[0].add_run()
                _set_run(run_t, text, font_name="Arial", size=10, color=INK)
        else:
            # Generic labeled list
            content_y = y + Inches(0.3)
            for label, text in labeled[:8]:
                tf, _ = _init_textbox(slide, content_x + Inches(0.2), content_y, content_w, Inches(0.4))
                p = tf.paragraphs[0]
                run_label = p.add_run()
                _set_run(run_label, f"{label}: ", font_name="Arial Black", size=11, bold=True, color=NAVY)
                run_text = p.add_run()
                _set_run(run_text, text, font_name="Arial", size=11, color=INK)
                content_y += Inches(0.45)

    # --- Step items as horizontal stepper (for Decision Path) ---
    elif step_items:
        content_y = y + Inches(0.5)
        n_steps = len(step_items)
        step_w = content_w / n_steps - Inches(0.15)

        # Horizontal rail line
        rail_y = content_y + Inches(0.35)
        _add_rect(slide, content_x + Inches(0.5), rail_y, content_w - Inches(1.0), Inches(0.03), fill_color=NEUTRAL)

        for i, (num, label, text) in enumerate(step_items):
            sx = content_x + i * (step_w + Inches(0.15))

            # Circle with number
            circle_size = Inches(0.55)
            cx = sx + (step_w - circle_size) / 2
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, content_y, circle_size, circle_size)
            circle.fill.solid()
            circle.fill.fore_color.rgb = NAVY if i < n_steps - 1 else RED
            circle.line.fill.background()
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_n = circle.text_frame.paragraphs[0].add_run()
            _set_run(run_n, num, font_name="Arial Black", size=14, bold=True, color=WHITE)

            # Step label
            tf_sl, _ = _init_textbox(slide, sx, content_y + Inches(0.7), step_w, Inches(0.35))
            tf_sl.paragraphs[0].alignment = PP_ALIGN.CENTER
            run_sl = tf_sl.paragraphs[0].add_run()
            _set_run(run_sl, label, font_name="Arial Black", size=10, bold=True, color=NAVY)

            # Step description
            tf_sd, _ = _init_textbox(slide, sx, content_y + Inches(1.1), step_w, Inches(0.8))
            tf_sd.paragraphs[0].alignment = PP_ALIGN.CENTER
            run_sd = tf_sd.paragraphs[0].add_run()
            _set_run(run_sd, text, font_name="Arial", size=9, color=MUTED)

    # --- Numbered items as bullet list ---
    elif numbered:
        content_y = y + Inches(0.3)
        for num, text in numbered[:8]:
            tf, _ = _init_textbox(slide, content_x + Inches(0.2), content_y, content_w, Inches(0.4))
            p = tf.paragraphs[0]
            run_num = p.add_run()
            _set_run(run_num, f"{num}.  ", font_name="Arial Black", size=12, bold=True, color=RED)
            run_text = p.add_run()
            _set_run(run_text, text, font_name="Arial", size=12, color=INK)
            content_y += Inches(0.45)

    # --- Bullet items ---
    elif body_items:
        content_y = y + Inches(0.3)
        for i, item in enumerate(body_items[:8]):
            tf, _ = _init_textbox(slide, content_x + Inches(0.2), content_y, content_w, Inches(0.4))
            p = tf.paragraphs[0]
            run_bullet = p.add_run()
            _set_run(run_bullet, "\u2022  ", font_name="Arial", size=12, color=RED)
            run_text = p.add_run()
            clean = re.sub(r'^\d+\.\s*', '', item).strip()
            _set_run(run_text, clean, font_name="Arial", size=12, color=INK)
            content_y += Inches(0.4)

    # --- Fallback: use narration as content if nothing else ---
    elif screen.get('narration') and len(screen['narration']) > 100:
        narration = screen['narration']
        # Split into sentences and show first few as key points
        sentences = [s.strip() + '.' for s in narration.split('.') if s.strip()]
        content_y = y + Inches(0.3)
        for sent in sentences[:5]:
            if len(sent) > 10:
                tf, _ = _init_textbox(slide, content_x + Inches(0.2), content_y, content_w, Inches(0.5))
                p = tf.paragraphs[0]
                run_bullet = p.add_run()
                _set_run(run_bullet, "\u2022  ", font_name="Arial", size=11, color=RED)
                run_text = p.add_run()
                _set_run(run_text, sent, font_name="Arial", size=11, color=INK)
                content_y += Inches(0.5)

    # Image placeholder if specified
    if has_image:
        img_x = Inches(8.0)
        img_y = Inches(2.8)
        img_w = Inches(4.5)
        img_h = Inches(3.0)
        _add_rounded_rect(slide, img_x, img_y, img_w, img_h, fill_color=NEUTRAL_2, border_color=NEUTRAL)
        tf_img, _ = _init_textbox(slide, img_x + Inches(0.5), img_y + Inches(1.2), img_w - Inches(1), Inches(0.5))
        p = tf_img.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        alt = screen.get('alt_text', screen.get('image', 'Image placeholder'))
        _set_run(run, f"[{alt}]", font_name="Arial", size=10, italic=True, color=MUTED)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_why_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 3: Why This Matters (s-why) — 2-column with stat block."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Healthcare relevance")
    y = _add_slide_title(slide, screen.get('title', 'Why This Matters'), y)

    # Left column: photo placeholder
    photo_x = CONTENT_PAD_X
    photo_y = Inches(2.8)
    photo_w = Inches(5.0)
    photo_h = Inches(3.5)
    _add_rounded_rect(slide, photo_x, photo_y, photo_w, photo_h, fill_color=NEUTRAL_2, border_color=NEUTRAL)
    tf_ph, _ = _init_textbox(slide, photo_x + Inches(1.2), photo_y + Inches(1.5), Inches(3), Inches(0.5))
    p = tf_ph.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, "[Photo Placeholder]", font_name="Arial", size=11, italic=True, color=MUTED)

    # Right column: stats + text
    right_x = Inches(6.5)
    right_w = Inches(6.0)

    # Stat block
    stat_bg = _add_rect(slide, right_x, Inches(2.8), Inches(3.0), Inches(1.2), fill_color=NAVY)
    tf_stat, _ = _init_textbox(slide, right_x + Inches(0.3), Inches(2.9), Inches(2.4), Inches(0.8))
    run_stat = tf_stat.paragraphs[0].add_run()
    _set_run(run_stat, "1 in 4", font_name="Arial Black", size=36, bold=True, color=WHITE)
    tf_stat_label, _ = _init_textbox(slide, right_x + Inches(0.3), Inches(3.6), Inches(2.4), Inches(0.3))
    run_sl = tf_stat_label.paragraphs[0].add_run()
    _set_run(run_sl, "Canadians have a disability", font_name="Arial", size=10, color=RGBColor(0xCC, 0xCC, 0xDD))

    # Text items
    body_items = [item for item in screen.get('text_items', [])
                  if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:'])]
    text_y = Inches(4.3)
    for item in body_items[:4]:
        tf, _ = _init_textbox(slide, right_x, text_y, right_w, Inches(0.35))
        p = tf.paragraphs[0]
        run = p.add_run()
        _set_run(run, f"\u2022  {item}", font_name="Arial", size=11, color=INK)
        text_y += Inches(0.35)

    # Key takeaway box (red left border)
    if screen.get('narration'):
        takeaway_y = Inches(5.8)
        _add_rect(slide, right_x, takeaway_y, Inches(0.06), Inches(0.7), fill_color=RED)
        _add_rect(slide, right_x + Inches(0.06), takeaway_y, right_w - Inches(0.06), Inches(0.7), fill_color=RED_BG)
        tf_tk, _ = _init_textbox(slide, right_x + Inches(0.25), takeaway_y + Inches(0.1), right_w - Inches(0.5), Inches(0.5))
        run_tk = tf_tk.paragraphs[0].add_run()
        _set_run(run_tk, "Key Takeaway", font_name="Arial Black", size=9, bold=True, color=RED)
        # Use first sentence of narration as takeaway summary
        narration = screen['narration']
        first_sentence = narration.split('.')[0] + '.' if '.' in narration else narration[:80]
        _add_paragraph(tf_tk, first_sentence, size=10, color=INK, space_before=4)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_impact_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Impact slide (s-impact) — two-column: photo left, accent bar + content right.
    Accent colour varies by topic: red for Missed Care, cobalt for Communication Gap, navy for Avoidance."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    title = screen.get('title', 'Impact')

    # Determine accent colour based on title keywords
    if 'Missed Care' in title:
        accent_color = RED
        accent_bg = RED_BG
    elif 'Communication' in title:
        accent_color = COBALT
        accent_bg = COBALT_BG
    else:
        # Avoidance or default
        accent_color = NAVY
        accent_bg = RGBColor(0xE8, 0xEB, 0xF2)  # light navy background

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Why this matters")
    y = _add_slide_title(slide, title, y)

    # --- LEFT COLUMN: photo placeholder ---
    photo_x = CONTENT_PAD_X
    photo_y = Inches(2.8)
    photo_w = Inches(5.0)
    photo_h = Inches(3.75)  # 540px equivalent
    _add_rounded_rect(slide, photo_x, photo_y, photo_w, photo_h,
                       fill_color=NEUTRAL_2, border_color=NEUTRAL)
    tf_ph, _ = _init_textbox(slide, photo_x + Inches(1.2), photo_y + Inches(1.6), Inches(3), Inches(0.5))
    p_ph = tf_ph.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    run_ph = p_ph.add_run()
    alt = screen.get('alt_text', '[Photo Placeholder]')
    _set_run(run_ph, f"[{alt}]", font_name="Arial", size=10, italic=True, color=MUTED)

    # --- RIGHT COLUMN: accent bar + H2 title + body text ---
    right_x = Inches(6.5)
    right_w = Inches(6.0)

    # Accent bar (vertical, left edge of right column)
    _add_rect(slide, right_x, Inches(2.8), Inches(0.06), Inches(3.75), fill_color=accent_color)

    # H2 title
    tf_h2, _ = _init_textbox(slide, right_x + Inches(0.3), Inches(2.9), right_w - Inches(0.5), Inches(0.5))
    run_h2 = tf_h2.paragraphs[0].add_run()
    # Strip "Impact: " prefix for cleaner heading
    display_title = re.sub(r'^Impact:\s*', '', title).strip()
    _set_run(run_h2, display_title, font_name="Arial Black", size=24, bold=True, color=accent_color)

    # Body text from text_items
    body_items = [item for item in screen.get('text_items', [])
                  if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:', 'H2:'])]
    text_y = Inches(3.6)
    for item in body_items[:3]:
        tf, _ = _init_textbox(slide, right_x + Inches(0.3), text_y, right_w - Inches(0.5), Inches(0.5))
        run = tf.paragraphs[0].add_run()
        _set_run(run, f"\u2022  {item}", font_name="Arial", size=11, color=INK)
        text_y += Inches(0.5)

    # Key takeaway / Indigenous context box
    box_y = Inches(5.3)
    is_indigenous = 'indigenous' in title.lower() or 'avoidance' in title.lower()
    if is_indigenous:
        # Earth-toned Indigenous context box
        earth_color = RGBColor(0x8B, 0x6F, 0x47)
        earth_bg = RGBColor(0xF5, 0xF0, 0xE6)
        _add_rect(slide, right_x + Inches(0.3), box_y, Inches(0.06), Inches(1.0), fill_color=earth_color)
        _add_rect(slide, right_x + Inches(0.36), box_y, right_w - Inches(0.86), Inches(1.0), fill_color=earth_bg)
        tf_ctx, _ = _init_textbox(slide, right_x + Inches(0.55), box_y + Inches(0.1),
                                   right_w - Inches(1.1), Inches(0.3))
        run_ctx_label = tf_ctx.paragraphs[0].add_run()
        _set_run(run_ctx_label, "Indigenous Context", font_name="Arial Black", size=9, bold=True, color=earth_color)
        # Use last bullet or narration excerpt
        indigenous_text = "Indigenous peoples face compounded barriers \u2014 systemic racism, geographic isolation, and culturally unsafe healthcare environments."
        _add_paragraph(tf_ctx, indigenous_text, size=10, color=INK, space_before=4)
    else:
        # Standard key takeaway box
        _add_rect(slide, right_x + Inches(0.3), box_y, Inches(0.06), Inches(0.8), fill_color=accent_color)
        _add_rect(slide, right_x + Inches(0.36), box_y, right_w - Inches(0.86), Inches(0.8), fill_color=accent_bg)
        tf_tk, _ = _init_textbox(slide, right_x + Inches(0.55), box_y + Inches(0.1),
                                  right_w - Inches(1.1), Inches(0.6))
        run_tk_label = tf_tk.paragraphs[0].add_run()
        _set_run(run_tk_label, "Key Takeaway", font_name="Arial Black", size=9, bold=True, color=accent_color)
        # Use first sentence of narration
        if screen.get('narration'):
            narration = screen['narration']
            first_sentence = narration.split('.')[0] + '.' if '.' in narration else narration[:80]
            _add_paragraph(tf_tk, first_sentence, size=10, color=INK, space_before=4)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_scenario_slides(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 7: Scenario (s-scenario) — setup + choice consequence slides."""
    slides = []

    # --- SETUP SLIDE ---
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Scenario")
    y = _add_slide_title(slide, screen.get('title', 'Scenario'), y)

    # Left: photo placeholder
    photo_x = CONTENT_PAD_X
    photo_y = Inches(2.8)
    _add_rounded_rect(slide, photo_x, photo_y, Inches(4.5), Inches(3.5),
                       fill_color=NEUTRAL_2, border_color=NEUTRAL)
    tf_ph, _ = _init_textbox(slide, photo_x + Inches(1.0), photo_y + Inches(1.5), Inches(2.5), Inches(0.5))
    p_ph = tf_ph.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    run_ph = p_ph.add_run()
    _set_run(run_ph, "[Character Photo]", font_name="Arial", size=11, italic=True, color=MUTED)

    # Right: scenario box with red left border
    box_x = Inches(5.8)
    box_y = Inches(2.8)
    box_w = Inches(6.8)
    _add_rect(slide, box_x, box_y, Inches(0.06), Inches(1.5), fill_color=RED)
    _add_rect(slide, box_x + Inches(0.06), box_y, box_w - Inches(0.06), Inches(1.5),
              fill_color=RED_BG)

    # Setup text
    tf_setup, _ = _init_textbox(slide, box_x + Inches(0.25), box_y + Inches(0.15),
                                 box_w - Inches(0.5), Inches(1.2))
    run_setup = tf_setup.paragraphs[0].add_run()
    _set_run(run_setup, screen.get('setup_text', 'Read the scenario below.'),
             font_name="Arial", size=11, color=INK)

    # "What do you do?" prompt
    tf_prompt, _ = _init_textbox(slide, box_x, Inches(4.6), box_w, Inches(0.35))
    run_prompt = tf_prompt.paragraphs[0].add_run()
    _set_run(run_prompt, "What do you do?", font_name="Arial Black", size=14, bold=True, color=COBALT)

    # Choice buttons (A/B/C)
    btn_y = Inches(5.1)
    choices = screen.get('choices', [])
    for i, choice in enumerate(choices):
        btn_x = box_x
        btn_w = box_w
        btn_h = Inches(0.45)

        # Button background
        _add_rounded_rect(slide, btn_x, btn_y, btn_w, btn_h, fill_color=WHITE, border_color=NEUTRAL)

        # Circle letter badge
        badge_size = Inches(0.3)
        badge_x = btn_x + Inches(0.1)
        badge_y_center = btn_y + (btn_h - badge_size) / 2
        badge = _add_rect(slide, badge_x, badge_y_center, badge_size, badge_size, fill_color=NAVY)
        # Make it a circle (use oval)
        badge_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y_center, badge_size, badge_size)
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = NAVY
        badge_shape.line.fill.background()
        # Letter on badge
        badge_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        run_letter = badge_shape.text_frame.paragraphs[0].add_run()
        _set_run(run_letter, choice['letter'], font_name="Arial Black", size=10, bold=True, color=WHITE)
        # Remove the rectangle badge (we used oval instead)
        badge._element.getparent().remove(badge._element)

        # Choice text
        tf_choice, _ = _init_textbox(slide, btn_x + Inches(0.5), btn_y, btn_w - Inches(0.6), btn_h)
        tf_choice.vertical_anchor = MSO_ANCHOR.MIDDLE
        run_c = tf_choice.paragraphs[0].add_run()
        _set_run(run_c, choice['text'], font_name="Arial", size=10, color=INK)

        btn_y += btn_h + Inches(0.1)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    slides.append(('setup', slide))

    # --- CONSEQUENCE SLIDES (one per choice) ---
    for ci, choice in enumerate(choices):
        c_slide = _add_blank_slide(prs)
        _add_topbar(c_slide, guide_num, section)

        quality = choice['quality']
        if 'Best' in quality:
            accent_color = CHARTREUSE
            accent_bg = CHARTREUSE_BG
        elif 'Poor' in quality:
            accent_color = RED
            accent_bg = RED_BG
        else:
            accent_color = COBALT
            accent_bg = COBALT_BG

        y = _add_eyebrow(c_slide, f"Choice {choice['letter']} \u00b7 {quality}")
        y = _add_slide_title(c_slide, f"Consequence: {choice['letter']}", y)

        # What you chose box
        chose_y = Inches(2.5)
        _add_rect(c_slide, CONTENT_PAD_X, chose_y, Inches(0.06), Inches(0.8), fill_color=accent_color)
        _add_rect(c_slide, CONTENT_PAD_X + Inches(0.06), chose_y, Inches(11.5), Inches(0.8),
                  fill_color=accent_bg)
        tf_chose, _ = _init_textbox(c_slide, CONTENT_PAD_X + Inches(0.25), chose_y + Inches(0.1),
                                     Inches(11), Inches(0.6))
        run_label = tf_chose.paragraphs[0].add_run()
        _set_run(run_label, "Your choice: ", font_name="Arial Black", size=10, bold=True, color=accent_color)
        run_text = tf_chose.paragraphs[0].add_run()
        _set_run(run_text, choice['text'], font_name="Arial", size=10, color=INK)

        # Consequence
        cons_y = Inches(3.6)
        tf_cons, _ = _init_textbox(c_slide, CONTENT_PAD_X, cons_y, Inches(11.5), Inches(1.5))
        run_rl = tf_cons.paragraphs[0].add_run()
        _set_run(run_rl, "Result: ", font_name="Arial Black", size=12, bold=True, color=accent_color)
        run_rt = tf_cons.paragraphs[0].add_run()
        _set_run(run_rt, choice['consequence'], font_name="Arial", size=12, color=INK)

        # Debrief on best choice
        if 'Best' in quality and screen.get('debrief'):
            db_y = Inches(5.0)
            _add_rect(c_slide, CONTENT_PAD_X, db_y, Inches(0.06), Inches(0.8), fill_color=COBALT)
            _add_rect(c_slide, CONTENT_PAD_X + Inches(0.06), db_y, Inches(11.5), Inches(0.8),
                      fill_color=COBALT_BG)
            tf_db, _ = _init_textbox(c_slide, CONTENT_PAD_X + Inches(0.25), db_y + Inches(0.1),
                                      Inches(11), Inches(0.6))
            run_db_label = tf_db.paragraphs[0].add_run()
            _set_run(run_db_label, "Key Takeaway: ", font_name="Arial Black", size=10, bold=True, color=COBALT)
            run_db = tf_db.paragraphs[0].add_run()
            _set_run(run_db, screen['debrief'], font_name="Arial", size=10, color=INK)

        # Principle
        if screen.get('principle') and 'Best' in quality:
            pr_y = Inches(6.1)
            tf_pr, _ = _init_textbox(c_slide, CONTENT_PAD_X, pr_y, Inches(11.5), Inches(0.3))
            run_pr = tf_pr.paragraphs[0].add_run()
            _set_run(run_pr, f"Principle: {screen['principle']}", font_name="Arial", size=9, italic=True, color=COBALT)

        _add_footer(c_slide, slide_idx + ci + 1, total_slides)
        slides.append((f'choice_{choice["letter"]}', c_slide))

    return slides


def create_question_slides(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 8: Knowledge Check (s-kc) — question + correct + incorrect feedback."""
    slides = []

    questions = screen.get('questions', [])
    if not questions:
        return [('content', create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides))]

    section = _extract_section_name(screen)

    for qi, q in enumerate(questions):
        q_num = qi + 1
        q_label = f"Screen {screen['number']}"
        if len(questions) > 1:
            q_label += f" - Q{q_num}"

        # --- QUESTION SLIDE ---
        q_slide = _add_blank_slide(prs)
        _add_topbar(q_slide, guide_num, section)

        section_num = screen.get('number', '').split('.')[0]
        y = _add_eyebrow(q_slide, f"Section {section_num} \u00b7 Knowledge check")
        y = _add_slide_title(q_slide, "Knowledge Check", y)

        # Left column: question + options
        q_x = CONTENT_PAD_X
        q_w = Inches(6.5)

        # Question stem
        stem_y = Inches(2.5)
        tf_stem, _ = _init_textbox(q_slide, q_x, stem_y, q_w, Inches(0.8))
        run_stem = tf_stem.paragraphs[0].add_run()
        _set_run(run_stem, q['stem'], font_name="Arial", size=14, bold=True, color=NAVY)

        # Options
        opt_y = Inches(3.5)
        for opt in q.get('options', []):
            opt_h = Inches(0.5)
            # Option button background
            _add_rounded_rect(q_slide, q_x, opt_y, q_w, opt_h, fill_color=WHITE, border_color=NEUTRAL)

            # Letter badge (circle)
            badge_size = Inches(0.3)
            badge_x = q_x + Inches(0.12)
            badge_y = opt_y + (opt_h - badge_size) / 2
            badge = q_slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y, badge_size, badge_size)
            badge.fill.solid()
            badge.fill.fore_color.rgb = NAVY
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = badge.text_frame.paragraphs[0].add_run()
            _set_run(run_l, opt['letter'].upper(), font_name="Arial Black", size=10, bold=True, color=WHITE)

            # Option text
            tf_opt, _ = _init_textbox(q_slide, q_x + Inches(0.55), opt_y, q_w - Inches(0.7), opt_h)
            tf_opt.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_o = tf_opt.paragraphs[0].add_run()
            _set_run(run_o, opt['text'], font_name="Arial", size=11, color=INK)

            opt_y += opt_h + Inches(0.12)

        # Right column: feedback panel placeholder
        panel_x = Inches(7.5)
        panel_w = Inches(5.2)
        panel_y = Inches(2.5)
        panel_h = Inches(3.5)
        _add_rounded_rect(q_slide, panel_x, panel_y, panel_w, panel_h,
                           fill_color=COBALT_BG, border_color=COBALT)
        tf_panel, _ = _init_textbox(q_slide, panel_x + Inches(0.3), panel_y + Inches(0.3),
                                     panel_w - Inches(0.6), Inches(0.4))
        run_panel = tf_panel.paragraphs[0].add_run()
        _set_run(run_panel, "FEEDBACK", font_name="Arial Black", size=10, bold=True, color=COBALT)
        tf_panel2, _ = _init_textbox(q_slide, panel_x + Inches(0.3), panel_y + Inches(0.8),
                                      panel_w - Inches(0.6), Inches(2.0))
        run_panel2 = tf_panel2.paragraphs[0].add_run()
        _set_run(run_panel2, "Select an answer to see feedback.",
                 font_name="Arial", size=11, italic=True, color=MUTED)

        _add_footer(q_slide, slide_idx + qi * 3, total_slides)
        _add_notes(q_slide, screen)
        slides.append(('question', q_slide))

        # --- CORRECT FEEDBACK SLIDE ---
        correct_opt = next((o for o in q.get('options', []) if o.get('correct')), None)
        c_slide = _add_blank_slide(prs)
        _add_topbar(c_slide, guide_num, section)

        y = _add_eyebrow(c_slide, f"Section {section_num} \u00b7 Knowledge check")
        y = _add_slide_title(c_slide, "Knowledge Check", y)

        # Left: question + highlighted correct answer
        tf_stem2, _ = _init_textbox(c_slide, q_x, Inches(2.5), q_w, Inches(0.8))
        run_s2 = tf_stem2.paragraphs[0].add_run()
        _set_run(run_s2, q['stem'], font_name="Arial", size=14, bold=True, color=NAVY)

        opt_y2 = Inches(3.5)
        for opt in q.get('options', []):
            opt_h = Inches(0.5)
            is_correct = opt.get('correct', False)
            bg = CHARTREUSE_BG if is_correct else WHITE
            border = CHARTREUSE if is_correct else NEUTRAL

            _add_rounded_rect(c_slide, q_x, opt_y2, q_w, opt_h, fill_color=bg, border_color=border)

            badge = c_slide.shapes.add_shape(MSO_SHAPE.OVAL, q_x + Inches(0.12),
                                              opt_y2 + (opt_h - Inches(0.3)) / 2,
                                              Inches(0.3), Inches(0.3))
            badge.fill.solid()
            badge.fill.fore_color.rgb = CHARTREUSE if is_correct else NAVY
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = badge.text_frame.paragraphs[0].add_run()
            symbol = "\u2713" if is_correct else opt['letter'].upper()
            _set_run(run_l, symbol, font_name="Arial Black", size=10, bold=True, color=WHITE)

            tf_o, _ = _init_textbox(c_slide, q_x + Inches(0.55), opt_y2, q_w - Inches(0.7), opt_h)
            tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_o = tf_o.paragraphs[0].add_run()
            _set_run(run_o, opt['text'], font_name="Arial", size=11,
                     color=CHARTREUSE if is_correct else MUTED)
            opt_y2 += opt_h + Inches(0.12)

        # Right: correct feedback panel
        _add_rounded_rect(c_slide, panel_x, panel_y, panel_w, panel_h,
                           fill_color=CHARTREUSE_BG, border_color=CHARTREUSE)
        tf_fb_title, _ = _init_textbox(c_slide, panel_x + Inches(0.3), panel_y + Inches(0.3),
                                        panel_w - Inches(0.6), Inches(0.5))
        run_fb_t = tf_fb_title.paragraphs[0].add_run()
        _set_run(run_fb_t, "\u2713  Correct!", font_name="Arial Black", size=16, bold=True, color=CHARTREUSE)

        feedback = screen.get('feedback_correct', 'Well done! You selected the correct answer.')
        tf_fb, _ = _init_textbox(c_slide, panel_x + Inches(0.3), panel_y + Inches(1.0),
                                  panel_w - Inches(0.6), Inches(2.0))
        run_fb = tf_fb.paragraphs[0].add_run()
        _set_run(run_fb, feedback, font_name="Arial", size=11, color=INK)

        if correct_opt:
            _add_paragraph(tf_fb, f"Answer: {correct_opt['letter']}) {correct_opt['text']}",
                          font_name="Arial", size=10, bold=True, color=CHARTREUSE, space_before=12)

        _add_footer(c_slide, slide_idx + qi * 3 + 1, total_slides)
        slides.append(('correct', c_slide))

        # --- INCORRECT FEEDBACK SLIDE ---
        i_slide = _add_blank_slide(prs)
        _add_topbar(i_slide, guide_num, section)

        y = _add_eyebrow(i_slide, f"Section {section_num} \u00b7 Knowledge check")
        y = _add_slide_title(i_slide, "Knowledge Check", y)

        # Left: question + highlighted wrong
        tf_stem3, _ = _init_textbox(i_slide, q_x, Inches(2.5), q_w, Inches(0.8))
        run_s3 = tf_stem3.paragraphs[0].add_run()
        _set_run(run_s3, q['stem'], font_name="Arial", size=14, bold=True, color=NAVY)

        opt_y3 = Inches(3.5)
        for opt in q.get('options', []):
            opt_h = Inches(0.5)
            is_correct = opt.get('correct', False)
            # On incorrect slide: show first non-correct option as "wrong"
            # The designer will assign this in Storyline; just mark correct in green
            bg = CHARTREUSE_BG if is_correct else WHITE
            border = CHARTREUSE if is_correct else NEUTRAL

            _add_rounded_rect(i_slide, q_x, opt_y3, q_w, opt_h, fill_color=bg, border_color=border)

            badge = i_slide.shapes.add_shape(MSO_SHAPE.OVAL, q_x + Inches(0.12),
                                              opt_y3 + (opt_h - Inches(0.3)) / 2,
                                              Inches(0.3), Inches(0.3))
            badge.fill.solid()
            badge.fill.fore_color.rgb = CHARTREUSE if is_correct else NAVY
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = badge.text_frame.paragraphs[0].add_run()
            _set_run(run_l, opt['letter'].upper(), font_name="Arial Black", size=10, bold=True, color=WHITE)

            tf_o, _ = _init_textbox(i_slide, q_x + Inches(0.55), opt_y3, q_w - Inches(0.7), opt_h)
            tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_o = tf_o.paragraphs[0].add_run()
            _set_run(run_o, opt['text'], font_name="Arial", size=11, color=INK)
            opt_y3 += opt_h + Inches(0.12)

        # Right: incorrect feedback panel
        _add_rounded_rect(i_slide, panel_x, panel_y, panel_w, panel_h,
                           fill_color=RED_BG, border_color=RED)
        tf_fb_title2, _ = _init_textbox(i_slide, panel_x + Inches(0.3), panel_y + Inches(0.3),
                                         panel_w - Inches(0.6), Inches(0.5))
        run_fb_t2 = tf_fb_title2.paragraphs[0].add_run()
        _set_run(run_fb_t2, "\u2717  Not quite.", font_name="Arial Black", size=16, bold=True, color=RED)

        feedback_inc = screen.get('feedback_incorrect', 'Review the content and try again.')
        tf_fb2, _ = _init_textbox(i_slide, panel_x + Inches(0.3), panel_y + Inches(1.0),
                                   panel_w - Inches(0.6), Inches(2.0))
        run_fb2 = tf_fb2.paragraphs[0].add_run()
        _set_run(run_fb2, feedback_inc, font_name="Arial", size=11, color=INK)

        if correct_opt:
            _add_paragraph(tf_fb2, f"The correct answer is: {correct_opt['letter']}) {correct_opt['text']}",
                          font_name="Arial", size=10, bold=True, color=CHARTREUSE, space_before=12)

        _add_footer(i_slide, slide_idx + qi * 3 + 2, total_slides)
        slides.append(('incorrect', i_slide))

    return slides


def create_tips_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 9: Inclusive Practice Tip (s-tip) — chartreuse accents."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Inclusive practice")
    y = _add_slide_title(slide, screen.get('title', 'Inclusive Practice Tip'), y)

    # Tip card
    card_x = CONTENT_PAD_X
    card_y = Inches(2.6)
    card_w = Inches(12.0)
    card_h = Inches(3.8)

    # Chartreuse left border
    _add_rect(slide, card_x, card_y, Inches(0.06), card_h, fill_color=CHARTREUSE)
    # Card background
    _add_rect(slide, card_x + Inches(0.06), card_y, card_w - Inches(0.06), card_h,
              fill_color=CHARTREUSE_BG, border_color=NEUTRAL, border_width=0.5)

    # Get tips from labeled_items
    tips = screen.get('labeled_items', [])
    if not tips:
        tips = [(f"Tip {i+1}", item) for i, item in enumerate(screen.get('text_items', []))
                if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:'])]

    inner_x = card_x + Inches(0.4)
    inner_y = card_y + Inches(0.3)
    inner_w = card_w - Inches(0.8)

    for i, (label, text) in enumerate(tips[:5]):
        # Tip label
        tf_title, _ = _init_textbox(slide, inner_x, inner_y, inner_w, Inches(0.25))
        run_t = tf_title.paragraphs[0].add_run()
        _set_run(run_t, label, font_name="Arial Black", size=10, bold=True, color=CHARTREUSE)

        # Tip text
        tf_content, _ = _init_textbox(slide, inner_x, inner_y + Inches(0.25), inner_w, Inches(0.4))
        run_c = tf_content.paragraphs[0].add_run()
        _set_run(run_c, text, font_name="Arial", size=10, color=INK)

        inner_y += Inches(0.7)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_map_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 10: MAP Activity (s-map) — lilac accents, 3 input fields."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 My action planning")
    y = _add_slide_title(slide, screen.get('title', 'My Accessibility Plan (MAP)'), y)

    # Subtitle
    y = _add_lede(slide, "Reflect on what you have learned. This is private and not submitted.", y)

    # 3 input fields from labeled_items (Stop/Start/Continue) or defaults
    labeled = screen.get('labeled_items', [])
    defaults = [
        ("STOP", "One thing I will stop doing"),
        ("START", "One thing I will start doing"),
        ("CONTINUE", "One thing I will continue doing"),
    ]
    fields = [(l.upper(), t.rstrip('*').strip()) for l, t in labeled[:3]] if labeled else defaults

    field_x = CONTENT_PAD_X
    field_w = Inches(12.0)
    field_h = Inches(1.1)
    field_y = Inches(3.2)

    for label, prompt in fields:
        # Field background
        _add_rect(slide, field_x, field_y, field_w, field_h, fill_color=LILAC_BG, border_color=LILAC, border_width=0.5)
        # Lilac left accent
        _add_rect(slide, field_x, field_y, Inches(0.06), field_h, fill_color=LILAC)

        # Label
        tf_label, _ = _init_textbox(slide, field_x + Inches(0.25), field_y + Inches(0.1), Inches(2), Inches(0.3))
        run_l = tf_label.paragraphs[0].add_run()
        _set_run(run_l, label, font_name="Arial Black", size=10, bold=True, color=LILAC)

        # Prompt text
        tf_prompt, _ = _init_textbox(slide, field_x + Inches(0.25), field_y + Inches(0.4),
                                      field_w - Inches(0.5), Inches(0.6))
        run_p = tf_prompt.paragraphs[0].add_run()
        _set_run(run_p, prompt, font_name="Arial", size=10, italic=True, color=MUTED)

        field_y += field_h + Inches(0.15)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_summary_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 11: Key Takeaways (s-take) — 2x2 numbered takeaway cards."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Review")
    y = _add_slide_title(slide, "Key Takeaways", y)

    # Prefer numbered_items, fallback to text_items
    items = [text for _, text in screen.get('numbered_items', [])]
    if not items:
        for item in screen.get('text_items', []):
            clean = re.sub(r'^\d+\.\s*', '', item)
            if clean and 'Audio' not in clean and 'Captions' not in clean:
                items.append(clean)

    # 2x2 grid of takeaway cards
    card_w = Inches(5.5)
    card_h = Inches(1.5)
    start_x = CONTENT_PAD_X
    start_y = Inches(2.8)
    gap_x = Inches(0.5)
    gap_y = Inches(0.35)

    for i, item in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y_card = start_y + row * (card_h + gap_y)

        # Card
        _add_rect(slide, x, y_card, card_w, card_h, fill_color=WHITE, border_color=NEUTRAL, border_width=0.5)
        _add_rect(slide, x, y_card, card_w, Inches(0.04), fill_color=COBALT)

        # Number
        tf_num, _ = _init_textbox(slide, x + Inches(0.15), y_card + Inches(0.12), Inches(0.5), Inches(0.5))
        run_num = tf_num.paragraphs[0].add_run()
        _set_run(run_num, str(i + 1), font_name="Arial Black", size=28, bold=True, color=COBALT)

        # Text
        tf_item, _ = _init_textbox(slide, x + Inches(0.7), y_card + Inches(0.12),
                                    card_w - Inches(0.85), card_h - Inches(0.2))
        run_item = tf_item.paragraphs[0].add_run()
        _set_run(run_item, item, font_name="Arial", size=11, color=INK)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_completion_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Final completion / resources slide — navy background with resources."""
    slide = _add_blank_slide(prs)

    # Navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Lilac accent line at top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), fill_color=LILAC)

    # UHN logo
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), CONTENT_PAD_X, Inches(0.5), height=Inches(0.5))

    # Completion title
    tf_title, _ = _init_textbox(slide, CONTENT_PAD_X, Inches(1.5), Inches(12), Inches(1.0))
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    subtitle = GUIDE_SUBTITLES.get(guide_num, "")
    _set_run(run, f"Guide {guide_num:02d} Complete", font_name="Arial Black", size=40, bold=True, color=WHITE)

    tf_sub, _ = _init_textbox(slide, CONTENT_PAD_X, Inches(2.6), Inches(12), Inches(0.5))
    p2 = tf_sub.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    _set_run(run2, subtitle, font_name="Arial", size=18, color=LILAC)

    # Resources list
    items = [i for i in screen.get('text_items', [])
             if 'Audio' not in i and 'Captions' not in i and 'badge' not in i.lower()]
    if items:
        res_y = Inches(3.5)
        tf_res_title, _ = _init_textbox(slide, Inches(3.0), res_y, Inches(7.333), Inches(0.4))
        p_rt = tf_res_title.paragraphs[0]
        p_rt.alignment = PP_ALIGN.CENTER
        run_rt = p_rt.add_run()
        _set_run(run_rt, "RESOURCES", font_name="Arial Black", size=10, bold=True, color=LILAC)

        res_y += Inches(0.5)
        for item in items[:6]:
            tf_r, _ = _init_textbox(slide, Inches(2.0), res_y, Inches(9.333), Inches(0.35))
            p_r = tf_r.paragraphs[0]
            p_r.alignment = PP_ALIGN.CENTER
            run_r = p_r.add_run()
            _set_run(run_r, item, font_name="Arial", size=11, color=RGBColor(0xCC, 0xCC, 0xDD))
            res_y += Inches(0.35)

    # "Next Guide" button
    btn_w = Inches(2.667)
    btn_x = (SLIDE_W - btn_w) / 2
    btn = _add_rounded_rect(slide, btn_x, Inches(6.0), btn_w, Inches(0.55), fill_color=RED)
    tf_btn, _ = _init_textbox(slide, btn_x, Inches(6.0), btn_w, Inches(0.55))
    tf_btn.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_btn = tf_btn.paragraphs[0]
    p_btn.alignment = PP_ALIGN.CENTER
    run_btn = p_btn.add_run()
    next_label = "NEXT GUIDE" if guide_num < 17 else "COMPLETE SERIES"
    _set_run(run_btn, next_label, font_name="Arial Black", size=12, bold=True, color=WHITE)

    # Bottom accent
    _add_rect(slide, Inches(0), FOOTER_Y + Inches(0.3), SLIDE_W, Inches(0.05), fill_color=LILAC)

    _add_notes(slide, screen)
    return slide


def create_reflection_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Reflection slide — cobalt accents."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Reflection")
    y = _add_slide_title(slide, screen.get('title', 'Reflection'), y)

    items = screen.get('text_items', [])
    items_clean = [i for i in items if 'Audio' not in i and 'Captions' not in i and 'keyboard' not in i]

    content_y = Inches(2.8)
    for item in items_clean[:6]:
        # Reflection prompt with cobalt left border
        _add_rect(slide, CONTENT_PAD_X, content_y, Inches(0.06), Inches(0.6), fill_color=COBALT)
        _add_rect(slide, CONTENT_PAD_X + Inches(0.06), content_y, Inches(11.5), Inches(0.6),
                  fill_color=COBALT_BG)
        tf, _ = _init_textbox(slide, CONTENT_PAD_X + Inches(0.25), content_y + Inches(0.1),
                               Inches(11), Inches(0.4))
        run = tf.paragraphs[0].add_run()
        is_private = 'private' in item.lower() or 'not submitted' in item.lower()
        _set_run(run, item, font_name="Arial", size=12, italic=is_private, color=INK if not is_private else MUTED)
        content_y += Inches(0.75)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_podcast_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 12: Listen and Reflect / Podcast (s-podcast) — cobalt accents."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Audio learning")
    # Use eyebrow color override to cobalt
    y = _add_slide_title(slide, "Listen and Reflect", y)
    y = _add_lede(slide, "A short patient-advisor reflection on what good accessibility actually feels like. Captions and a full transcript are always available.", y)

    # Audio player placeholder
    player_x = CONTENT_PAD_X
    player_y = Inches(3.0)
    player_w = Inches(5.5)
    player_h = Inches(2.5)
    _add_rounded_rect(slide, player_x, player_y, player_w, player_h,
                       fill_color=COBALT_BG, border_color=COBALT)

    # Episode label
    tf_ep, _ = _init_textbox(slide, player_x + Inches(0.3), player_y + Inches(0.2), player_w - Inches(0.6), Inches(0.3))
    run_ep = tf_ep.paragraphs[0].add_run()
    _set_run(run_ep, f"Episode {guide_num:02d} of 18 \u00b7 {GUIDE_SUBTITLES.get(guide_num, '')}", font_name="Arial", size=10, color=COBALT)

    # Podcast title
    tf_pt, _ = _init_textbox(slide, player_x + Inches(0.3), player_y + Inches(0.6), player_w - Inches(0.6), Inches(0.5))
    run_pt = tf_pt.paragraphs[0].add_run()
    _set_run(run_pt, "Accessibility in Everyday Care", font_name="Arial Black", size=16, bold=True, color=NAVY)

    # Progress bar placeholder
    bar_y = player_y + Inches(1.4)
    _add_rect(slide, player_x + Inches(0.3), bar_y, player_w - Inches(0.6), Inches(0.06), fill_color=NEUTRAL)
    _add_rect(slide, player_x + Inches(0.3), bar_y, Inches(1.5), Inches(0.06), fill_color=COBALT)

    # Controls text
    tf_ctrl, _ = _init_textbox(slide, player_x + Inches(0.3), bar_y + Inches(0.15), player_w - Inches(0.6), Inches(0.3))
    run_ctrl = tf_ctrl.paragraphs[0].add_run()
    _set_run(run_ctrl, "\u25B6  0:53 / 2:30    CC  Captions on    1.0\u00d7", font_name="Arial", size=9, color=MUTED)

    # Key listening points (right side)
    kp_x = Inches(7.0)
    kp_w = Inches(5.5)
    tf_kp_title, _ = _init_textbox(slide, kp_x, Inches(3.0), kp_w, Inches(0.3))
    run_kpt = tf_kp_title.paragraphs[0].add_run()
    _set_run(run_kpt, "Key listening points", font_name="Arial Black", size=12, bold=True, color=NAVY)

    # Extract listening points from text_items or labeled_items
    points = screen.get('text_items', [])
    points = [p for p in points if p.startswith('00:') or p.startswith('01:') or p.startswith('02:')]
    if not points:
        points = ["00:18 \u2014 Why a tone of voice at intake matters more than the form.",
                   "01:02 \u2014 How \"ask, then confirm\" prevented a wrong assumption.",
                   "01:54 \u2014 One sentence every staff member can use on a busy shift."]

    kp_y = Inches(3.5)
    for point in points[:4]:
        _add_rect(slide, kp_x, kp_y, kp_w, Inches(0.55), fill_color=NEUTRAL_2, border_color=NEUTRAL, border_width=0.5)
        tf_kp, _ = _init_textbox(slide, kp_x + Inches(0.15), kp_y + Inches(0.08), kp_w - Inches(0.3), Inches(0.4))
        run_kp = tf_kp.paragraphs[0].add_run()
        _set_run(run_kp, point, font_name="Arial", size=10, color=INK)
        kp_y += Inches(0.65)

    # Reflection prompt at bottom
    ref_y = Inches(6.0)
    _add_rect(slide, CONTENT_PAD_X, ref_y, Inches(12.0), Inches(0.06), fill_color=COBALT)
    tf_ref, _ = _init_textbox(slide, CONTENT_PAD_X, ref_y + Inches(0.15), Inches(12.0), Inches(0.4))
    run_ref = tf_ref.paragraphs[0].add_run()
    _set_run(run_ref, "Reflection: Think of one patient interaction this week. Which moment would have shifted if you had paused to ask, before you acted?",
             font_name="Arial", size=10, italic=True, color=MUTED)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_progress_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 14: Series Progress Map (s-progress) — 3-stage journey."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Series progress \u00b7 Stages unlock as you complete")
    y = _add_slide_title(slide, "Accessibility First \u00b7 A three-stage journey", y)
    y = _add_lede(slide, "Guides unlock by stage. Foundations open everything else.", y)

    # 3 stages
    stages = [
        ("Stage 1", "Foundations", "Guides 01\u201304 \u00b7 Required first", RED, [
            ("01", "Foundations", "Complete" if guide_num >= 1 else ""),
            ("02", "Perceptions & Attitudes", ""),
            ("03", "Vision Disabilities", ""),
            ("04", "Hearing & Communication", ""),
        ]),
        ("Stage 2", "Understanding", "Guides 05\u201309 \u00b7 Builds on Foundations", COBALT, [
            ("05", "Physical & Mobility", ""),
            ("06", "Mental Health", ""),
            ("07", "Intellectual & Dev.", ""),
            ("08", "Non-Visible", ""),
            ("09", "Aging & Intersect.", ""),
        ]),
        ("Stage 3", "Applied Practice", "Guides 10\u201318 \u00b7 Unlock after Stages 1+2", CHARTREUSE, [
            ("10", "Confidence & Respect", ""),
            ("11", "Service Animals", ""),
            ("12", "Support Persons", ""),
        ]),
    ]

    stage_y = Inches(3.0)
    for stage_label, stage_name, stage_desc, color, guides in stages:
        # Stage header
        _add_rect(slide, CONTENT_PAD_X, stage_y, Inches(12.0), Inches(0.4), fill_color=color)
        tf_sh, _ = _init_textbox(slide, CONTENT_PAD_X + Inches(0.15), stage_y + Inches(0.05), Inches(4), Inches(0.3))
        run_sh = tf_sh.paragraphs[0].add_run()
        _set_run(run_sh, f"{stage_label}: {stage_name}", font_name="Arial Black", size=11, bold=True, color=WHITE)

        tf_sd, _ = _init_textbox(slide, Inches(5.0), stage_y + Inches(0.05), Inches(7), Inches(0.3))
        run_sd = tf_sd.paragraphs[0].add_run()
        _set_run(run_sd, stage_desc, font_name="Arial", size=9, color=WHITE)

        # Guide items
        item_y = stage_y + Inches(0.45)
        for g_num, g_name, g_status in guides[:4]:
            tf_g, _ = _init_textbox(slide, CONTENT_PAD_X + Inches(0.3), item_y, Inches(11.5), Inches(0.25))
            p = tf_g.paragraphs[0]
            run_n = p.add_run()
            _set_run(run_n, f"{g_num}  ", font_name="Arial Black", size=9, bold=True, color=color)
            run_gn = p.add_run()
            _set_run(run_gn, g_name, font_name="Arial", size=9, color=INK)
            if g_status:
                run_s = p.add_run()
                _set_run(run_s, f"  \u2022 {g_status}", font_name="Arial", size=8, color=CHARTREUSE)
            item_y += Inches(0.25)

        stage_y = item_y + Inches(0.15)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


# ---------------------------------------------------------------------------
# Main generation
# ---------------------------------------------------------------------------

def _count_total_slides(screens):
    """Pre-count total slides for progress dots."""
    count = 1  # title slide
    for screen in screens:
        stype = screen['type']
        if stype in ('scenario', 'decision_tree'):
            count += 1 + len(screen.get('choices', []))
        elif stype == 'question':
            questions = screen.get('questions', [])
            count += max(len(questions), 1) * 3
        else:
            count += 1
    return count


def generate_ppt(guide_num: int, dry_run: bool = False):
    """Generate the full PPT from the master storyboard."""
    guide_folder = BUILD_OUTPUT / GUIDE_FOLDERS[guide_num]
    storyboard_path = guide_folder / "01-master-storyboard" / f"MASTER-STORYBOARD-GUIDE-{guide_num:02d}.md"

    if not storyboard_path.exists():
        print(f"Error: Storyboard not found: {storyboard_path}")
        sys.exit(1)

    print(f"Parsing storyboard: {storyboard_path.name}")
    data = parse_storyboard(storyboard_path)
    info = data['info']
    screens = data['screens']

    print(f"Found {len(screens)} screens")
    print(f"Course: {info.get('title', 'Unknown')}")
    print(f"Seat time: {info.get('seat_time', 'N/A')}")
    print()

    total_slides = _count_total_slides(screens)

    # Plan slides
    slide_plan = []
    slide_plan.append(('title', 'Course Title Slide', None))

    for screen in screens:
        stype = screen['type']
        label = f"Screen {screen['number']} -- {screen['title']}"

        if stype in ('scenario', 'decision_tree'):
            slide_plan.append(('scenario_setup', f"{label} (Setup)", screen))
            for choice in screen.get('choices', []):
                slide_plan.append(('scenario_choice', f"  -> Choice {choice['letter']} ({choice['quality']})", screen))
        elif stype == 'question':
            questions = screen.get('questions', [])
            for qi, q in enumerate(questions):
                q_label = f"{label}"
                if len(questions) > 1:
                    q_label += f" Q{qi+1}"
                slide_plan.append(('question', f"{q_label} (Question)", screen))
                slide_plan.append(('correct', f"  -> Correct Feedback", screen))
                slide_plan.append(('incorrect', f"  -> Incorrect Feedback", screen))
        else:
            slide_plan.append((stype, label, screen))

    print(f"Slide plan ({len(slide_plan)} slides):")
    print("-" * 60)
    for i, (stype, label, _) in enumerate(slide_plan, 1):
        print(f"  {i:3d}. [{stype:16s}] {label}")
    print("-" * 60)
    print()

    if dry_run:
        print("Dry run -- no PPT generated.")
        return

    # Generate PPT
    print("Generating PPT (mockup-matched, shapes-from-scratch)...")

    # Build from scratch — no template, pure mockup design
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_idx = 0

    # 1. Title slide
    create_title_slide(prs, info, guide_num, total_slides)
    print("  [ok] Title slide")
    slide_idx += 1

    # 2. Screen slides
    for screen in screens:
        stype = screen['type']
        label = f"Screen {screen['number']} -- {screen['title']}"

        if stype == 'welcome':
            create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (welcome)")
            slide_idx += 1

        elif stype == 'objectives':
            create_objectives_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (objectives)")
            slide_idx += 1

        elif stype == 'impact':
            create_impact_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (impact)")
            slide_idx += 1

        elif stype == 'scenario':
            scenario_slides = create_scenario_slides(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (scenario: {len(scenario_slides)} slides)")
            slide_idx += len(scenario_slides)

        elif stype == 'question':
            q_slides = create_question_slides(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (questions: {len(q_slides)} slides)")
            slide_idx += len(q_slides)

        elif stype == 'action_planning':
            create_map_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (MAP)")
            slide_idx += 1

        elif stype == 'reflection':
            create_reflection_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (reflection)")
            slide_idx += 1

        elif stype == 'tips':
            create_tips_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (tip)")
            slide_idx += 1

        elif stype == 'summary':
            create_summary_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (takeaways)")
            slide_idx += 1

        elif stype == 'decision_tree':
            dt_slides = create_scenario_slides(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (decision tree: {len(dt_slides)} slides)")
            slide_idx += len(dt_slides)

        elif stype == 'podcast':
            create_podcast_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (podcast)")
            slide_idx += 1

        elif stype == 'progress':
            create_progress_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (progress)")
            slide_idx += 1

        elif stype == 'completion':
            create_completion_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (completion)")
            slide_idx += 1

        else:
            create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (content)")
            slide_idx += 1

    # Save
    output_dir = guide_folder / "02-ppt-storyboard"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"UHN_Guide{guide_num:02d}_Storyboard.pptx"
    prs.save(str(output_path))

    actual_slides = len(prs.slides)
    print(f"\n[ok] PPT saved: {output_path}")
    print(f"  Total slides: {actual_slides}")
    print(f"  Dimensions: 13.333\" x 7.5\" (16:9 widescreen)")
    print(f"  Size: {output_path.stat().st_size / 1024:.0f} KB")


def main():
    parser = argparse.ArgumentParser(description="Generate PPT storyboard for UHN Accessibility First courses")
    parser.add_argument("--guide", type=int, required=True, help="Guide number (1-18)")
    parser.add_argument("--dry-run", action="store_true", help="Preview slide plan without generating")
    args = parser.parse_args()

    if args.guide not in GUIDE_FOLDERS:
        print(f"Error: Guide {args.guide} not found. Valid: 1-18")
        sys.exit(1)

    generate_ppt(args.guide, dry_run=args.dry_run)


if __name__ == "__main__":
    main()
