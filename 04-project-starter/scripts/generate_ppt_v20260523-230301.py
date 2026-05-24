#!/usr/bin/env python3
"""
PPT Storyboard Generator for UHN Accessibility First Course Series
Converts Master Storyboard markdown into branded PowerPoint for Storyline 360 import.

Builds every slide from scratch using python-pptx shapes on a blank layout,
matching the HTML mockup design system (1920x1080 / 16:9).

Each screen becomes one or more slides:
- Content screens → 1 slide
- Question screens → question slide + correct feedback slide + incorrect feedback slide
- Branching scenarios → setup slide + consequence slides per choice (A/B/C) + debrief slide

Usage:
    python generate_ppt.py --guide 1
    python generate_ppt.py --guide 1 --dry-run    # Preview slide plan without generating

Requirements:
    pip install python-pptx Pillow
"""

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("Error: 'python-pptx' package not installed. Run: pip install python-pptx")
    sys.exit(1)

BASE_DIR = Path(__file__).resolve().parent.parent.parent
BUILD_OUTPUT = BASE_DIR / "05-build-output"
TEMPLATE_PATH = BASE_DIR / "02-branding-and-style" / "uhn-templates" / "UHN_Presentation_Template_Primary (1).pptx"
LOGO_PATH = BASE_DIR / "02-branding-and-style" / "logos" / "uhn_logo.png"

GUIDE_FOLDERS = {
    1: "01-Foundations-of-Disability-Inclusion-and-Accessible-Design",
    2: "02-Perceptions-Attitudes-and-Barriers",
    3: "03-Vision-Disabilities",
    4: "04-Sensory-Hearing-and-Communication-Disabilities",
    5: "05-Physical-Disabilities-and-Mobility",
    6: "06-Mental-Health-Disabilities",
    7: "07-Intellectual-Developmental-and-Learning-Disabilities",
    8: "08-Non-Visible-Disabilities",
    9: "09-Aging-Disability-and-Intersectionality",
    10: "10-Engaging-with-Confidence-and-Respect",
    11: "11-Service-Animals-Guide-Dogs-and-Non-Service-Animals",
    12: "12-Support-Persons",
    13: "13-Assistive-Devices",
    14: "14-Communication-and-Information-Accessibility",
    15: "15-Neurodiversity-and-Sensory-Regulation",
    16: "16-Trauma-Informed-Accessibility",
    17: "17-Accessibility-in-Crisis-Situations-and-De-escalation",
    18: "18-Indigenous-Peoples-and-Accessibility",
}

# Guide subtitle lookup
GUIDE_SUBTITLES = {
    1: "Foundations",
    2: "Perceptions & Attitudes",
    3: "Vision Disabilities",
    4: "Hearing & Communication",
    5: "Physical & Mobility",
    6: "Mental Health",
    7: "Intellectual & Developmental",
    8: "Non-Visible Disabilities",
    9: "Aging & Intersectionality",
    10: "Confidence & Respect",
    11: "Service Animals",
    12: "Support Persons",
    13: "Assistive Devices",
    14: "Communication Access",
    15: "Neurodiversity",
    16: "Trauma-Informed",
    17: "Crisis & De-escalation",
    18: "Indigenous Peoples",
}

# ---------------------------------------------------------------------------
# Design system colours (from HTML mockup CSS variables)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x19, 0x28, 0x58)
NAVY_DEEP = RGBColor(0x0F, 0x1A, 0x3D)
RED = RGBColor(0xC0, 0x23, 0x3B)
RED_BG = RGBColor(0xFD, 0xF2, 0xF3)
COBALT = RGBColor(0x24, 0x5B, 0xAA)
COBALT_BG = RGBColor(0xEA, 0xF1, 0xFA)
LILAC = RGBColor(0xC4, 0x8A, 0xBD)
LILAC_BG = RGBColor(0xF7, 0xED, 0xF5)
CHARTREUSE = RGBColor(0x74, 0xAE, 0x54)
CHARTREUSE_BG = RGBColor(0xEE, 0xF6, 0xE7)
NEUTRAL = RGBColor(0xE4, 0xE4, 0xE4)
NEUTRAL_2 = RGBColor(0xF4, 0xF4, 0xF4)
INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Common chrome dimensions
TOPBAR_H = Inches(0.583)        # 84px
CONTENT_PAD_X = Inches(0.667)   # 96px
CONTENT_TOP = Inches(0.75)      # below topbar
FOOTER_Y = Inches(6.95)
FOOTER_H = Inches(0.417)        # ~60px


# ---------------------------------------------------------------------------
# Helper: text formatting
# ---------------------------------------------------------------------------

def _set_run(run, text, font_name="Arial", size=14, bold=False, italic=False, color=None):
    """Apply text formatting to a run."""
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def _add_paragraph(tf, text, font_name="Arial", size=14, bold=False, italic=False,
                   color=None, alignment=None, space_before=0, space_after=0):
    """Add a paragraph with a single run to a text frame."""
    p = tf.add_paragraph()
    if alignment:
        p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    run = p.add_run()
    _set_run(run, text, font_name=font_name, size=size, bold=bold, italic=italic, color=color)
    return p


def _init_textbox(slide, left, top, width, height, word_wrap=True):
    """Add a textbox and return the text frame (first paragraph cleared)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    # Clear the default empty paragraph
    tf.paragraphs[0].clear()
    return tf, txBox


def _add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # no border by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()  # transparent
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = Pt(border_width)
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
    return shape


# ---------------------------------------------------------------------------
# Common chrome: topbar + footer
# ---------------------------------------------------------------------------

def _add_topbar(slide, guide_num, section_name=""):
    """Add the navy topbar with UHN logo + breadcrumbs."""
    # Navy background bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, TOPBAR_H, fill_color=NAVY)

    # UHN logo
    logo_h = Inches(0.333)  # 48px
    logo_top = Inches((0.583 - 0.333) / 2)  # centered in topbar
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.333), logo_top, height=logo_h)

    # Separator line (thin vertical)
    if LOGO_PATH.exists():
        sep_x = Inches(1.5)
    else:
        sep_x = Inches(0.333)
    _add_rect(slide, sep_x, Inches(0.125), Inches(0.014), Inches(0.333), fill_color=RGBColor(0x4A, 0x5A, 0x8A))

    # "Accessibility First" label
    tf, _ = _init_textbox(slide, sep_x + Inches(0.15), Inches(0.05), Inches(2.5), TOPBAR_H)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    _set_run(run, "Accessibility First", font_name="Arial Black", size=11, bold=True, color=WHITE)

    # Breadcrumbs on the right
    subtitle = GUIDE_SUBTITLES.get(guide_num, "")
    breadcrumb = f"Guide {guide_num:02d} \u00b7 {subtitle}"
    if section_name:
        breadcrumb += f"  >  {section_name}"
    bc_tf, _ = _init_textbox(slide, Inches(8.5), Inches(0.05), Inches(4.5), TOPBAR_H)
    bc_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = bc_tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    _set_run(run, breadcrumb, font_name="Arial", size=9, color=RGBColor(0xAA, 0xBB, 0xDD))


def _add_footer(slide, slide_index, total_slides):
    """Add footer with ACCESSIBILITY FIRST tag + progress dots."""
    # Red dot
    _add_rect(slide, Inches(0.667), FOOTER_Y + Inches(0.08), Inches(0.083), Inches(0.083), fill_color=RED)

    # "ACCESSIBILITY FIRST" tag
    tf, _ = _init_textbox(slide, Inches(0.83), FOOTER_Y, Inches(2.5), Inches(0.35))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    _set_run(run, "ACCESSIBILITY FIRST", font_name="Arial Black", size=7, bold=True, color=NAVY)

    # Progress dots on the right
    dot_w = Inches(0.18)
    dot_h = Inches(0.055)
    dot_gap = Inches(0.07)
    # Max 20 dots to avoid overcrowding
    max_dots = min(total_slides, 20)
    step = max(1, total_slides // max_dots)
    dots_to_show = (total_slides + step - 1) // step
    start_x = Inches(13.333) - CONTENT_PAD_X - (dots_to_show * (dot_w + dot_gap))
    dot_y = FOOTER_Y + Inches(0.07)

    for i in range(dots_to_show):
        actual_idx = i * step
        x = start_x + i * (dot_w + dot_gap)
        if actual_idx < slide_index:
            color = NAVY
        elif actual_idx == slide_index:
            color = RED
        else:
            color = NEUTRAL
        _add_rect(slide, x, dot_y, dot_w, dot_h, fill_color=color)


def _add_eyebrow(slide, text, y=None):
    """Add an eyebrow label (red, uppercase, letter-spaced)."""
    if y is None:
        y = CONTENT_TOP
    tf, _ = _init_textbox(slide, CONTENT_PAD_X, y, Inches(10), Inches(0.35))
    run = tf.paragraphs[0].add_run()
    # Simulate letter-spacing with spaces
    spaced = "  ".join(text.upper())
    _set_run(run, spaced, font_name="Arial Black", size=9, bold=True, color=RED)
    return y + Inches(0.35)


def _add_slide_title(slide, text, y=None):
    """Add the main slide H1 title."""
    if y is None:
        y = CONTENT_TOP + Inches(0.35)
    tf, _ = _init_textbox(slide, CONTENT_PAD_X, y, Inches(11), Inches(0.7))
    run = tf.paragraphs[0].add_run()
    _set_run(run, text, font_name="Arial Black", size=36, bold=True, color=NAVY)
    return y + Inches(0.7)


def _add_lede(slide, text, y=None):
    """Add a lede / subtitle paragraph."""
    if y is None:
        y = CONTENT_TOP + Inches(1.1)
    tf, _ = _init_textbox(slide, CONTENT_PAD_X, y, Inches(10), Inches(0.6))
    run = tf.paragraphs[0].add_run()
    _set_run(run, text, font_name="Arial", size=15, color=MUTED)
    return y + Inches(0.6)


def _add_notes(slide, screen):
    """Add narration + audio info to slide notes."""
    notes_text = ""
    if screen.get('narration'):
        notes_text = f"NARRATION:\n{screen['narration']}"
    if screen.get('audio'):
        notes_text += f"\n\nAudio: {screen['audio']}"
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def _get_blank_layout(prs):
    """Get the blank slide layout (or the closest one with fewest placeholders)."""
    # Try to find a layout named "Blank"
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            return layout
    # Fallback: use the layout with fewest placeholders
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _add_blank_slide(prs):
    """Add a blank slide with white background."""
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    # Clear any placeholder shapes inherited from layout
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    # Set white background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


# ---------------------------------------------------------------------------
# Storyboard parser (UNCHANGED)
# ---------------------------------------------------------------------------

def parse_storyboard(filepath: Path) -> dict:
    """Parse the master storyboard markdown into structured data."""
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    # Extract header info
    info = {}
    title_match = re.search(r'^# (.+)', content, re.MULTILINE)
    if title_match:
        info['title'] = title_match.group(1).strip()

    seat_match = re.search(r'Seat Time\s*\|\s*(.+?)(?:\s*\|)', content)
    if seat_match:
        info['seat_time'] = seat_match.group(1).strip()

    screens_match = re.search(r'Screens\s*\|\s*(\d+)', content)
    if screens_match:
        info['total_screens'] = int(screens_match.group(1))

    clo_match = re.search(r'\*\*CLOs?:\*\*\s*(.+?)$', content, re.MULTILINE)
    if clo_match:
        info['clos'] = clo_match.group(1).strip()

    # Parse individual screens
    screen_sections = re.split(r'(?=^## Screen \d)', content, flags=re.MULTILINE)
    screens = []

    for section in screen_sections:
        if not section.strip().startswith('## Screen'):
            continue

        screen = parse_screen(section)
        if screen:
            screens.append(screen)

    return {'info': info, 'screens': screens}


def flatten_table_column(section: str, col_idx: int) -> str:
    """Extract all content from a specific column of a markdown table.
    col_idx: 0=Step, 1=Activities, 2=Design Guide"""
    lines = []
    for line in section.split('\n'):
        line = line.strip()
        if not line.startswith('|') or line.startswith('|---') or line.startswith('| Step'):
            continue
        cells = [c.strip() for c in line.split('|')]
        # cells[0] is empty (before first |), cells[-1] is empty (after last |)
        if len(cells) > col_idx + 1:
            cell = cells[col_idx + 1].strip()
            if cell:
                lines.append(cell)
    return '\n'.join(lines)


def parse_screen(section: str) -> dict:
    """Parse a single screen section into structured data."""
    screen = {}

    # Screen number and title
    header_match = re.search(r'## Screen ([\d.]+)\s*[—–-]\s*(.+)', section)
    if header_match:
        screen['number'] = header_match.group(1).strip()
        screen['title'] = header_match.group(2).strip()

    # Flatten the Activities column (col 1) for easier parsing
    activities = flatten_table_column(section, 1)
    design_col = flatten_table_column(section, 2)

    # Determine screen type based on title and content
    title = screen.get('title', '')
    screen['type'] = 'content'  # default

    if 'Welcome' in title:
        screen['type'] = 'welcome'
    elif 'Learning Objectives' in title or 'Objectives' in title:
        screen['type'] = 'objectives'
    elif 'Scenario' in title and ('Branching' in activities or 'Choice A' in activities):
        screen['type'] = 'scenario'
    elif 'Knowledge Check' in title or 'Assessment' in title or 'Question' in activities[:50]:
        screen['type'] = 'question'
    elif 'Reflection' in title:
        screen['type'] = 'reflection'
    elif 'MAP' in title or 'Action Planning' in title:
        screen['type'] = 'action_planning'
    elif 'Completion' in title or 'Resources' in title:
        screen['type'] = 'completion'
    elif 'Key Takeaways' in title or 'Summary' in title:
        screen['type'] = 'summary'
    elif 'Inclusive Practice' in title or 'Tips' in title:
        screen['type'] = 'tips'

    # Extract narration — find the narration text in activities column
    narration_match = re.search(
        r'\*\*Narration\s*\([^)]*\):\*\*\s*\n?(.*?)(?:\*\*Audio|\*\*Refs|\*\*Audio:|\Z)',
        activities, re.DOTALL
    )
    if narration_match:
        narration = narration_match.group(1).strip()
        # Clean up any remaining markdown
        narration = re.sub(r'\*\*[^*]+\*\*', '', narration).strip()
        if narration:
            screen['narration'] = narration

    # Extract setup text for scenarios (text before "What do you do?")
    if screen['type'] == 'scenario':
        setup_match = re.search(
            r'(?:Setup Text:|Scenario — Branching)\s*\n?(.*?)(?:\*\*What do you do\?\*\*|What do you do\?)',
            activities, re.DOTALL
        )
        if not setup_match:
            # Try getting the long text block from activities
            lines = activities.split('\n')
            setup_lines = []
            for line in lines:
                if 'Choice A' in line or 'What do you do' in line:
                    break
                clean = line.strip().lstrip('*').rstrip('*').strip()
                if clean and not any(x in clean for x in ['Setup Text:', 'Scenario', 'Interaction:', 'Time:', 'SME:', 'Screen 1.']):
                    setup_lines.append(clean)
            if setup_lines:
                screen['setup_text'] = ' '.join(setup_lines)
        else:
            text = setup_match.group(1).strip()
            # Remove markdown formatting artifacts
            text = re.sub(r'\*\*[^*]+:\*\*', '', text).strip()
            text = re.sub(r'\n+', ' ', text).strip()
            screen['setup_text'] = text

    # Extract choices for scenarios
    choices = []
    choice_pattern = re.finditer(
        r'\*\*Choice ([A-C])\s*\(([^)]+)\):\*\*\s*(.*?)→\s*\*(.*?)\*',
        activities, re.DOTALL
    )
    for m in choice_pattern:
        letter, quality, text, consequence = m.group(1), m.group(2), m.group(3), m.group(4)
        choices.append({
            'letter': letter.strip(),
            'quality': quality.strip(),
            'text': re.sub(r'\n+', ' ', text).strip(),
            'consequence': re.sub(r'\n+', ' ', consequence).strip(),
        })
    if choices:
        screen['choices'] = choices

    # Extract debrief
    debrief_match = re.search(r'\*\*Debrief:\*\*\s*(.*?)(?:\*\*Audio|\*\*Principle|\Z)', activities, re.DOTALL)
    if debrief_match:
        screen['debrief'] = re.sub(r'\n+', ' ', debrief_match.group(1)).strip()

    # Extract questions (for knowledge check screens)
    questions = []
    q_iter = re.finditer(
        r'\*\*Question \d+\s*\(([^)]+)\):\*\*\s*\n?(.*?)(?=\*\*Question \d+|\*\*Audio|\Z)',
        activities, re.DOTALL
    )
    for m in q_iter:
        q_type = m.group(1).strip()
        q_body = m.group(2).strip()
        q = {'type': q_type}

        # Parse question text and options from the body
        lines = [l.strip() for l in q_body.split('\n') if l.strip()]
        # First line is the stem
        stem_lines = []
        options = []
        for line in lines:
            opt_match = re.match(r'^([a-d])\)\s*(.+?)(\s*\u2713)?$', line)
            if opt_match:
                options.append({
                    'letter': opt_match.group(1),
                    'text': opt_match.group(2).strip(),
                    'correct': bool(opt_match.group(3)),
                })
            elif not options:  # still building the stem
                stem_lines.append(line)

        q['stem'] = ' '.join(stem_lines)
        q['options'] = options
        questions.append(q)

    if questions:
        screen['questions'] = questions

    # Extract feedback from design column
    correct_match = re.search(r'Correct:\s*"([^"]+)"', design_col)
    if correct_match:
        screen['feedback_correct'] = correct_match.group(1).strip()
    incorrect_match = re.search(r'Incorrect:\s*"([^"]+)"', design_col)
    if incorrect_match:
        screen['feedback_incorrect'] = incorrect_match.group(1).strip()

    # Extract on-screen text items from activities column (• bullets)
    text_items = re.findall(r'\u2022\s*(.+?)$', activities, re.MULTILINE)
    screen['text_items'] = [t.strip() for t in text_items if t.strip()]

    # Extract numbered list items (1. xxx, 2. xxx)
    numbered_items = re.findall(r'^(\d+)\.\s+(.+?)$', activities, re.MULTILINE)
    screen['numbered_items'] = [(num, text.strip()) for num, text in numbered_items if text.strip()]

    # Extract bold-label items (**Label:** text) — for models, frameworks, tips, steps
    skip_labels = {'Narration', 'Audio', 'Setup Text', 'On-Screen Text', 'Debrief',
                   'Refs', 'Prompt', 'Completion message', 'Next in series',
                   'Resources', 'Screen', 'What do you do?'}
    labeled_items = re.findall(r'\*\*([^*]+?):\*\*\s*(.+?)$', activities, re.MULTILINE)
    screen['labeled_items'] = [
        (label.strip(), text.strip()) for label, text in labeled_items
        if label.strip() not in skip_labels
        and not label.startswith('Question')
        and not label.startswith('Choice')
        and not label.startswith('Narration')
        and not re.match(r'Screen \d', label)
    ]

    # Extract bold step items (**1. Label** — text) — for decision path
    step_items = re.findall(r'\*\*(\d+)\.\s*([^*]+)\*\*\s*[—–-]\s*(.+?)$', activities, re.MULTILINE)
    screen['step_items'] = [(num, label.strip(), text.strip()) for num, label, text in step_items]

    # Extract image filename from design column
    img_match = re.search(r'\*\*Image:\*\*\s*(.+?)$', design_col, re.MULTILINE)
    if img_match:
        screen['image'] = img_match.group(1).strip()

    # Extract alt text from design column
    alt_match = re.search(r'Alt:\s*"(.+?)"', design_col)
    if alt_match:
        screen['alt_text'] = alt_match.group(1).strip()

    # Extract audio file
    audio_match = re.search(r'\*\*Audio:\*\*\s*(voiceover_[\d.]+\.mp3)', activities)
    if audio_match:
        screen['audio'] = audio_match.group(1)

    # Extract principle from design column
    principle_match = re.search(r'\*\*Principle:\*\*\s*(.+?)$', design_col, re.MULTILINE)
    if principle_match:
        screen['principle'] = principle_match.group(1).strip()

    return screen


# ---------------------------------------------------------------------------
# Slide builders — mockup-matched, all from scratch
# ---------------------------------------------------------------------------

def _extract_section_name(screen):
    """Try to derive section name from screen number/title."""
    title = screen.get('title', '')
    # Strip markdown artifacts
    title = re.sub(r'\*+', '', title).strip()
    return title


def create_title_slide(prs, info, guide_num, total_slides):
    """Slide 1: Title slide (s-title) — Navy background, 2-column."""
    slide = _add_blank_slide(prs)

    # Navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Lilac accent line at top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), fill_color=LILAC)

    # --- LEFT COLUMN (x=0.667 to ~6.5) ---
    left_x = CONTENT_PAD_X

    # UHN logo (top-left)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left_x, Inches(0.5), height=Inches(0.5))

    # Series tag: "Guide XX of 17 · Foundations"
    subtitle = GUIDE_SUBTITLES.get(guide_num, "")
    series_text = f"Guide {guide_num:02d} of 17  \u00b7  {subtitle}"
    tf, _ = _init_textbox(slide, left_x, Inches(1.3), Inches(5.5), Inches(0.4))
    run = tf.paragraphs[0].add_run()
    _set_run(run, series_text, font_name="Arial", size=14, color=LILAC)

    # Main title (H1)
    raw_title = info.get('title', 'Accessibility First')
    # Strip "Master Storyboard — Guide XX: " prefix if present
    clean_title = re.sub(r'Master Storyboard\s*[—–-]\s*Guide\s*\d+:\s*', '', raw_title).strip()
    tf, _ = _init_textbox(slide, left_x, Inches(1.9), Inches(5.5), Inches(1.8))
    run = tf.paragraphs[0].add_run()
    _set_run(run, clean_title, font_name="Arial Black", size=42, bold=True, color=WHITE)

    # Subtitle
    tf, _ = _init_textbox(slide, left_x, Inches(3.8), Inches(5.5), Inches(0.6))
    run = tf.paragraphs[0].add_run()
    _set_run(run, "Accessibility First Course Series", font_name="Arial", size=16, color=RGBColor(0xCC, 0xCC, 0xDD))

    # Meta row
    seat_time = info.get('seat_time', '15-20 min')
    total_screens = info.get('total_screens', '')
    meta_items = [
        f"Duration: {seat_time}",
        "Format: eLearning (Storyline 360)",
        "Audience: All UHN staff",
    ]
    meta_y = Inches(4.7)
    for item in meta_items:
        tf, _ = _init_textbox(slide, left_x, meta_y, Inches(5), Inches(0.3))
        run = tf.paragraphs[0].add_run()
        _set_run(run, item, font_name="Arial", size=10, color=MUTED)
        meta_y += Inches(0.3)

    # --- RIGHT COLUMN (x=7.0 to ~12.5) ---
    right_x = Inches(7.0)
    right_w = Inches(5.667)

    # Photo placeholder
    photo_rect = _add_rounded_rect(slide, right_x, Inches(0.8), right_w, Inches(4.5),
                                    fill_color=RGBColor(0x2A, 0x3A, 0x6A))
    tf2, _ = _init_textbox(slide, right_x + Inches(1.5), Inches(2.7), Inches(3), Inches(0.5))
    run2 = tf2.paragraphs[0].add_run()
    _set_run(run2, "[Hero Photo Placeholder]", font_name="Arial", size=12, color=MUTED)

    # BEGIN button
    btn = _add_rounded_rect(slide, right_x + Inches(1.5), Inches(5.8), Inches(2.667), Inches(0.55),
                             fill_color=RED)
    tf_btn, _ = _init_textbox(slide, right_x + Inches(1.5), Inches(5.8), Inches(2.667), Inches(0.55))
    tf_btn.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_btn.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run_btn = p.add_run()
    _set_run(run_btn, "BEGIN", font_name="Arial Black", size=14, bold=True, color=WHITE)

    # Footer (no progress dots on title)
    _add_rect(slide, Inches(0), FOOTER_Y, SLIDE_W, Inches(0.05), fill_color=LILAC)

    return slide


def create_objectives_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 2: Learning Objectives (s-obj) — 2x2 grid of objective cards."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    # Eyebrow
    section_num = screen.get('number', '2').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 What you will learn")

    # Title
    y = _add_slide_title(slide, "Learning Objectives", y)

    # Objective cards (2x2 grid)
    objectives = [item for item in screen.get('text_items', [])
                  if any(item.startswith(f"{n}.") for n in range(1, 10))]
    if not objectives:
        objectives = [item for item in screen.get('text_items', [])
                      if 'CLO' not in item and 'Audio' not in item]

    card_w = Inches(5.5)
    card_h = Inches(1.3)
    start_x = CONTENT_PAD_X
    start_y = Inches(2.8)
    gap_x = Inches(0.5)
    gap_y = Inches(0.35)

    for i, obj_text in enumerate(objectives[:4]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y_card = start_y + row * (card_h + gap_y)

        # Card background
        _add_rect(slide, x, y_card, card_w, card_h, fill_color=WHITE, border_color=NEUTRAL, border_width=0.5)
        # Navy top border accent
        _add_rect(slide, x, y_card, card_w, Inches(0.04), fill_color=NAVY)

        # Number
        tf_num, _ = _init_textbox(slide, x + Inches(0.15), y_card + Inches(0.12), Inches(0.5), Inches(0.5))
        run_num = tf_num.paragraphs[0].add_run()
        _set_run(run_num, str(i + 1), font_name="Arial Black", size=28, bold=True, color=NAVY)

        # Objective text (strip leading number)
        clean = re.sub(r'^\d+\.\s*', '', obj_text).strip()
        tf_obj, _ = _init_textbox(slide, x + Inches(0.7), y_card + Inches(0.12), card_w - Inches(0.85), card_h - Inches(0.2))
        run_obj = tf_obj.paragraphs[0].add_run()
        _set_run(run_obj, clean, font_name="Arial", size=11, color=INK)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides):
    """Generic content slide with topbar, eyebrow, title, and bullet list."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 {screen.get('title', '')}")
    y = _add_slide_title(slide, screen.get('title', 'Content'), y)

    # Body content
    body_items = [item for item in screen.get('text_items', [])
                  if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:'])]

    if body_items:
        content_y = y + Inches(0.3)
        for i, item in enumerate(body_items[:8]):
            tf, _ = _init_textbox(slide, CONTENT_PAD_X + Inches(0.2), content_y, Inches(10), Inches(0.4))
            # Bullet
            p = tf.paragraphs[0]
            run_bullet = p.add_run()
            _set_run(run_bullet, "\u2022  ", font_name="Arial", size=12, color=RED)
            run_text = p.add_run()
            clean = re.sub(r'^\d+\.\s*', '', item).strip()
            _set_run(run_text, clean, font_name="Arial", size=12, color=INK)
            content_y += Inches(0.4)

    # Image placeholder if specified
    if screen.get('image'):
        img_x = Inches(8.0)
        img_y = Inches(2.8)
        img_w = Inches(4.5)
        img_h = Inches(3.0)
        _add_rounded_rect(slide, img_x, img_y, img_w, img_h, fill_color=NEUTRAL_2, border_color=NEUTRAL)
        tf_img, _ = _init_textbox(slide, img_x + Inches(0.5), img_y + Inches(1.2), img_w - Inches(1), Inches(0.5))
        p = tf_img.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        alt = screen.get('alt_text', screen.get('image', 'Image placeholder'))
        _set_run(run, f"[{alt}]", font_name="Arial", size=10, italic=True, color=MUTED)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_why_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 3: Why This Matters (s-why) — 2-column with stat block."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Healthcare relevance")
    y = _add_slide_title(slide, screen.get('title', 'Why This Matters'), y)

    # Left column: photo placeholder
    photo_x = CONTENT_PAD_X
    photo_y = Inches(2.8)
    photo_w = Inches(5.0)
    photo_h = Inches(3.5)
    _add_rounded_rect(slide, photo_x, photo_y, photo_w, photo_h, fill_color=NEUTRAL_2, border_color=NEUTRAL)
    tf_ph, _ = _init_textbox(slide, photo_x + Inches(1.2), photo_y + Inches(1.5), Inches(3), Inches(0.5))
    p = tf_ph.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, "[Photo Placeholder]", font_name="Arial", size=11, italic=True, color=MUTED)

    # Right column: stats + text
    right_x = Inches(6.5)
    right_w = Inches(6.0)

    # Stat block
    stat_bg = _add_rect(slide, right_x, Inches(2.8), Inches(3.0), Inches(1.2), fill_color=NAVY)
    tf_stat, _ = _init_textbox(slide, right_x + Inches(0.3), Inches(2.9), Inches(2.4), Inches(0.8))
    run_stat = tf_stat.paragraphs[0].add_run()
    _set_run(run_stat, "1 in 4", font_name="Arial Black", size=36, bold=True, color=WHITE)
    tf_stat_label, _ = _init_textbox(slide, right_x + Inches(0.3), Inches(3.6), Inches(2.4), Inches(0.3))
    run_sl = tf_stat_label.paragraphs[0].add_run()
    _set_run(run_sl, "Canadians have a disability", font_name="Arial", size=10, color=RGBColor(0xCC, 0xCC, 0xDD))

    # Text items
    body_items = [item for item in screen.get('text_items', [])
                  if not any(skip in item for skip in ['Audio:', 'Captions:', 'Image:', 'SME:'])]
    text_y = Inches(4.3)
    for item in body_items[:4]:
        tf, _ = _init_textbox(slide, right_x, text_y, right_w, Inches(0.35))
        p = tf.paragraphs[0]
        run = p.add_run()
        _set_run(run, f"\u2022  {item}", font_name="Arial", size=11, color=INK)
        text_y += Inches(0.35)

    # Key takeaway box (red left border)
    if screen.get('narration'):
        takeaway_y = Inches(5.8)
        _add_rect(slide, right_x, takeaway_y, Inches(0.06), Inches(0.7), fill_color=RED)
        _add_rect(slide, right_x + Inches(0.06), takeaway_y, right_w - Inches(0.06), Inches(0.7), fill_color=RED_BG)
        tf_tk, _ = _init_textbox(slide, right_x + Inches(0.25), takeaway_y + Inches(0.1), right_w - Inches(0.5), Inches(0.5))
        run_tk = tf_tk.paragraphs[0].add_run()
        _set_run(run_tk, "Key Takeaway", font_name="Arial Black", size=9, bold=True, color=RED)
        # Use first sentence of narration as takeaway summary
        narration = screen['narration']
        first_sentence = narration.split('.')[0] + '.' if '.' in narration else narration[:80]
        _add_paragraph(tf_tk, first_sentence, size=10, color=INK, space_before=4)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_scenario_slides(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 7: Scenario (s-scenario) — setup + choice consequence slides."""
    slides = []

    # --- SETUP SLIDE ---
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Scenario")
    y = _add_slide_title(slide, screen.get('title', 'Scenario'), y)

    # Left: photo placeholder
    photo_x = CONTENT_PAD_X
    photo_y = Inches(2.8)
    _add_rounded_rect(slide, photo_x, photo_y, Inches(4.5), Inches(3.5),
                       fill_color=NEUTRAL_2, border_color=NEUTRAL)
    tf_ph, _ = _init_textbox(slide, photo_x + Inches(1.0), photo_y + Inches(1.5), Inches(2.5), Inches(0.5))
    p_ph = tf_ph.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    run_ph = p_ph.add_run()
    _set_run(run_ph, "[Character Photo]", font_name="Arial", size=11, italic=True, color=MUTED)

    # Right: scenario box with red left border
    box_x = Inches(5.8)
    box_y = Inches(2.8)
    box_w = Inches(6.8)
    _add_rect(slide, box_x, box_y, Inches(0.06), Inches(1.5), fill_color=RED)
    _add_rect(slide, box_x + Inches(0.06), box_y, box_w - Inches(0.06), Inches(1.5),
              fill_color=RED_BG)

    # Setup text
    tf_setup, _ = _init_textbox(slide, box_x + Inches(0.25), box_y + Inches(0.15),
                                 box_w - Inches(0.5), Inches(1.2))
    run_setup = tf_setup.paragraphs[0].add_run()
    _set_run(run_setup, screen.get('setup_text', 'Read the scenario below.'),
             font_name="Arial", size=11, color=INK)

    # "What do you do?" prompt
    tf_prompt, _ = _init_textbox(slide, box_x, Inches(4.6), box_w, Inches(0.35))
    run_prompt = tf_prompt.paragraphs[0].add_run()
    _set_run(run_prompt, "What do you do?", font_name="Arial Black", size=14, bold=True, color=COBALT)

    # Choice buttons (A/B/C)
    btn_y = Inches(5.1)
    choices = screen.get('choices', [])
    for i, choice in enumerate(choices):
        btn_x = box_x
        btn_w = box_w
        btn_h = Inches(0.45)

        # Button background
        _add_rounded_rect(slide, btn_x, btn_y, btn_w, btn_h, fill_color=WHITE, border_color=NEUTRAL)

        # Circle letter badge
        badge_size = Inches(0.3)
        badge_x = btn_x + Inches(0.1)
        badge_y_center = btn_y + (btn_h - badge_size) / 2
        badge = _add_rect(slide, badge_x, badge_y_center, badge_size, badge_size, fill_color=NAVY)
        # Make it a circle (use oval)
        badge_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y_center, badge_size, badge_size)
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = NAVY
        badge_shape.line.fill.background()
        # Letter on badge
        badge_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        run_letter = badge_shape.text_frame.paragraphs[0].add_run()
        _set_run(run_letter, choice['letter'], font_name="Arial Black", size=10, bold=True, color=WHITE)
        # Remove the rectangle badge (we used oval instead)
        badge._element.getparent().remove(badge._element)

        # Choice text
        tf_choice, _ = _init_textbox(slide, btn_x + Inches(0.5), btn_y, btn_w - Inches(0.6), btn_h)
        tf_choice.vertical_anchor = MSO_ANCHOR.MIDDLE
        run_c = tf_choice.paragraphs[0].add_run()
        _set_run(run_c, choice['text'], font_name="Arial", size=10, color=INK)

        btn_y += btn_h + Inches(0.1)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    slides.append(('setup', slide))

    # --- CONSEQUENCE SLIDES (one per choice) ---
    for ci, choice in enumerate(choices):
        c_slide = _add_blank_slide(prs)
        _add_topbar(c_slide, guide_num, section)

        quality = choice['quality']
        if 'Best' in quality:
            accent_color = CHARTREUSE
            accent_bg = CHARTREUSE_BG
        elif 'Poor' in quality:
            accent_color = RED
            accent_bg = RED_BG
        else:
            accent_color = COBALT
            accent_bg = COBALT_BG

        y = _add_eyebrow(c_slide, f"Choice {choice['letter']} \u00b7 {quality}")
        y = _add_slide_title(c_slide, f"Consequence: {choice['letter']}", y)

        # What you chose box
        chose_y = Inches(2.5)
        _add_rect(c_slide, CONTENT_PAD_X, chose_y, Inches(0.06), Inches(0.8), fill_color=accent_color)
        _add_rect(c_slide, CONTENT_PAD_X + Inches(0.06), chose_y, Inches(11.5), Inches(0.8),
                  fill_color=accent_bg)
        tf_chose, _ = _init_textbox(c_slide, CONTENT_PAD_X + Inches(0.25), chose_y + Inches(0.1),
                                     Inches(11), Inches(0.6))
        run_label = tf_chose.paragraphs[0].add_run()
        _set_run(run_label, "Your choice: ", font_name="Arial Black", size=10, bold=True, color=accent_color)
        run_text = tf_chose.paragraphs[0].add_run()
        _set_run(run_text, choice['text'], font_name="Arial", size=10, color=INK)

        # Consequence
        cons_y = Inches(3.6)
        tf_cons, _ = _init_textbox(c_slide, CONTENT_PAD_X, cons_y, Inches(11.5), Inches(1.5))
        run_rl = tf_cons.paragraphs[0].add_run()
        _set_run(run_rl, "Result: ", font_name="Arial Black", size=12, bold=True, color=accent_color)
        run_rt = tf_cons.paragraphs[0].add_run()
        _set_run(run_rt, choice['consequence'], font_name="Arial", size=12, color=INK)

        # Debrief on best choice
        if 'Best' in quality and screen.get('debrief'):
            db_y = Inches(5.0)
            _add_rect(c_slide, CONTENT_PAD_X, db_y, Inches(0.06), Inches(0.8), fill_color=COBALT)
            _add_rect(c_slide, CONTENT_PAD_X + Inches(0.06), db_y, Inches(11.5), Inches(0.8),
                      fill_color=COBALT_BG)
            tf_db, _ = _init_textbox(c_slide, CONTENT_PAD_X + Inches(0.25), db_y + Inches(0.1),
                                      Inches(11), Inches(0.6))
            run_db_label = tf_db.paragraphs[0].add_run()
            _set_run(run_db_label, "Key Takeaway: ", font_name="Arial Black", size=10, bold=True, color=COBALT)
            run_db = tf_db.paragraphs[0].add_run()
            _set_run(run_db, screen['debrief'], font_name="Arial", size=10, color=INK)

        # Principle
        if screen.get('principle') and 'Best' in quality:
            pr_y = Inches(6.1)
            tf_pr, _ = _init_textbox(c_slide, CONTENT_PAD_X, pr_y, Inches(11.5), Inches(0.3))
            run_pr = tf_pr.paragraphs[0].add_run()
            _set_run(run_pr, f"Principle: {screen['principle']}", font_name="Arial", size=9, italic=True, color=COBALT)

        _add_footer(c_slide, slide_idx + ci + 1, total_slides)
        slides.append((f'choice_{choice["letter"]}', c_slide))

    return slides


def create_question_slides(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 8: Knowledge Check (s-kc) — question + correct + incorrect feedback."""
    slides = []

    questions = screen.get('questions', [])
    if not questions:
        return [('content', create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides))]

    section = _extract_section_name(screen)

    for qi, q in enumerate(questions):
        q_num = qi + 1
        q_label = f"Screen {screen['number']}"
        if len(questions) > 1:
            q_label += f" - Q{q_num}"

        # --- QUESTION SLIDE ---
        q_slide = _add_blank_slide(prs)
        _add_topbar(q_slide, guide_num, section)

        section_num = screen.get('number', '').split('.')[0]
        y = _add_eyebrow(q_slide, f"Section {section_num} \u00b7 Knowledge check")
        y = _add_slide_title(q_slide, "Knowledge Check", y)

        # Left column: question + options
        q_x = CONTENT_PAD_X
        q_w = Inches(6.5)

        # Question stem
        stem_y = Inches(2.5)
        tf_stem, _ = _init_textbox(q_slide, q_x, stem_y, q_w, Inches(0.8))
        run_stem = tf_stem.paragraphs[0].add_run()
        _set_run(run_stem, q['stem'], font_name="Arial", size=14, bold=True, color=NAVY)

        # Options
        opt_y = Inches(3.5)
        for opt in q.get('options', []):
            opt_h = Inches(0.5)
            # Option button background
            _add_rounded_rect(q_slide, q_x, opt_y, q_w, opt_h, fill_color=WHITE, border_color=NEUTRAL)

            # Letter badge (circle)
            badge_size = Inches(0.3)
            badge_x = q_x + Inches(0.12)
            badge_y = opt_y + (opt_h - badge_size) / 2
            badge = q_slide.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y, badge_size, badge_size)
            badge.fill.solid()
            badge.fill.fore_color.rgb = NAVY
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = badge.text_frame.paragraphs[0].add_run()
            _set_run(run_l, opt['letter'].upper(), font_name="Arial Black", size=10, bold=True, color=WHITE)

            # Option text
            tf_opt, _ = _init_textbox(q_slide, q_x + Inches(0.55), opt_y, q_w - Inches(0.7), opt_h)
            tf_opt.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_o = tf_opt.paragraphs[0].add_run()
            _set_run(run_o, opt['text'], font_name="Arial", size=11, color=INK)

            opt_y += opt_h + Inches(0.12)

        # Right column: feedback panel placeholder
        panel_x = Inches(7.5)
        panel_w = Inches(5.2)
        panel_y = Inches(2.5)
        panel_h = Inches(3.5)
        _add_rounded_rect(q_slide, panel_x, panel_y, panel_w, panel_h,
                           fill_color=COBALT_BG, border_color=COBALT)
        tf_panel, _ = _init_textbox(q_slide, panel_x + Inches(0.3), panel_y + Inches(0.3),
                                     panel_w - Inches(0.6), Inches(0.4))
        run_panel = tf_panel.paragraphs[0].add_run()
        _set_run(run_panel, "FEEDBACK", font_name="Arial Black", size=10, bold=True, color=COBALT)
        tf_panel2, _ = _init_textbox(q_slide, panel_x + Inches(0.3), panel_y + Inches(0.8),
                                      panel_w - Inches(0.6), Inches(2.0))
        run_panel2 = tf_panel2.paragraphs[0].add_run()
        _set_run(run_panel2, "Select an answer to see feedback.",
                 font_name="Arial", size=11, italic=True, color=MUTED)

        _add_footer(q_slide, slide_idx + qi * 3, total_slides)
        _add_notes(q_slide, screen)
        slides.append(('question', q_slide))

        # --- CORRECT FEEDBACK SLIDE ---
        correct_opt = next((o for o in q.get('options', []) if o.get('correct')), None)
        c_slide = _add_blank_slide(prs)
        _add_topbar(c_slide, guide_num, section)

        y = _add_eyebrow(c_slide, f"Section {section_num} \u00b7 Knowledge check")
        y = _add_slide_title(c_slide, "Knowledge Check", y)

        # Left: question + highlighted correct answer
        tf_stem2, _ = _init_textbox(c_slide, q_x, Inches(2.5), q_w, Inches(0.8))
        run_s2 = tf_stem2.paragraphs[0].add_run()
        _set_run(run_s2, q['stem'], font_name="Arial", size=14, bold=True, color=NAVY)

        opt_y2 = Inches(3.5)
        for opt in q.get('options', []):
            opt_h = Inches(0.5)
            is_correct = opt.get('correct', False)
            bg = CHARTREUSE_BG if is_correct else WHITE
            border = CHARTREUSE if is_correct else NEUTRAL

            _add_rounded_rect(c_slide, q_x, opt_y2, q_w, opt_h, fill_color=bg, border_color=border)

            badge = c_slide.shapes.add_shape(MSO_SHAPE.OVAL, q_x + Inches(0.12),
                                              opt_y2 + (opt_h - Inches(0.3)) / 2,
                                              Inches(0.3), Inches(0.3))
            badge.fill.solid()
            badge.fill.fore_color.rgb = CHARTREUSE if is_correct else NAVY
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = badge.text_frame.paragraphs[0].add_run()
            symbol = "\u2713" if is_correct else opt['letter'].upper()
            _set_run(run_l, symbol, font_name="Arial Black", size=10, bold=True, color=WHITE)

            tf_o, _ = _init_textbox(c_slide, q_x + Inches(0.55), opt_y2, q_w - Inches(0.7), opt_h)
            tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_o = tf_o.paragraphs[0].add_run()
            _set_run(run_o, opt['text'], font_name="Arial", size=11,
                     color=CHARTREUSE if is_correct else MUTED)
            opt_y2 += opt_h + Inches(0.12)

        # Right: correct feedback panel
        _add_rounded_rect(c_slide, panel_x, panel_y, panel_w, panel_h,
                           fill_color=CHARTREUSE_BG, border_color=CHARTREUSE)
        tf_fb_title, _ = _init_textbox(c_slide, panel_x + Inches(0.3), panel_y + Inches(0.3),
                                        panel_w - Inches(0.6), Inches(0.5))
        run_fb_t = tf_fb_title.paragraphs[0].add_run()
        _set_run(run_fb_t, "\u2713  Correct!", font_name="Arial Black", size=16, bold=True, color=CHARTREUSE)

        feedback = screen.get('feedback_correct', 'Well done! You selected the correct answer.')
        tf_fb, _ = _init_textbox(c_slide, panel_x + Inches(0.3), panel_y + Inches(1.0),
                                  panel_w - Inches(0.6), Inches(2.0))
        run_fb = tf_fb.paragraphs[0].add_run()
        _set_run(run_fb, feedback, font_name="Arial", size=11, color=INK)

        if correct_opt:
            _add_paragraph(tf_fb, f"Answer: {correct_opt['letter']}) {correct_opt['text']}",
                          font_name="Arial", size=10, bold=True, color=CHARTREUSE, space_before=12)

        _add_footer(c_slide, slide_idx + qi * 3 + 1, total_slides)
        slides.append(('correct', c_slide))

        # --- INCORRECT FEEDBACK SLIDE ---
        i_slide = _add_blank_slide(prs)
        _add_topbar(i_slide, guide_num, section)

        y = _add_eyebrow(i_slide, f"Section {section_num} \u00b7 Knowledge check")
        y = _add_slide_title(i_slide, "Knowledge Check", y)

        # Left: question + highlighted wrong
        tf_stem3, _ = _init_textbox(i_slide, q_x, Inches(2.5), q_w, Inches(0.8))
        run_s3 = tf_stem3.paragraphs[0].add_run()
        _set_run(run_s3, q['stem'], font_name="Arial", size=14, bold=True, color=NAVY)

        opt_y3 = Inches(3.5)
        for opt in q.get('options', []):
            opt_h = Inches(0.5)
            is_correct = opt.get('correct', False)
            # On incorrect slide: show first non-correct option as "wrong"
            # The designer will assign this in Storyline; just mark correct in green
            bg = CHARTREUSE_BG if is_correct else WHITE
            border = CHARTREUSE if is_correct else NEUTRAL

            _add_rounded_rect(i_slide, q_x, opt_y3, q_w, opt_h, fill_color=bg, border_color=border)

            badge = i_slide.shapes.add_shape(MSO_SHAPE.OVAL, q_x + Inches(0.12),
                                              opt_y3 + (opt_h - Inches(0.3)) / 2,
                                              Inches(0.3), Inches(0.3))
            badge.fill.solid()
            badge.fill.fore_color.rgb = CHARTREUSE if is_correct else NAVY
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = badge.text_frame.paragraphs[0].add_run()
            _set_run(run_l, opt['letter'].upper(), font_name="Arial Black", size=10, bold=True, color=WHITE)

            tf_o, _ = _init_textbox(i_slide, q_x + Inches(0.55), opt_y3, q_w - Inches(0.7), opt_h)
            tf_o.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_o = tf_o.paragraphs[0].add_run()
            _set_run(run_o, opt['text'], font_name="Arial", size=11, color=INK)
            opt_y3 += opt_h + Inches(0.12)

        # Right: incorrect feedback panel
        _add_rounded_rect(i_slide, panel_x, panel_y, panel_w, panel_h,
                           fill_color=RED_BG, border_color=RED)
        tf_fb_title2, _ = _init_textbox(i_slide, panel_x + Inches(0.3), panel_y + Inches(0.3),
                                         panel_w - Inches(0.6), Inches(0.5))
        run_fb_t2 = tf_fb_title2.paragraphs[0].add_run()
        _set_run(run_fb_t2, "\u2717  Not quite.", font_name="Arial Black", size=16, bold=True, color=RED)

        feedback_inc = screen.get('feedback_incorrect', 'Review the content and try again.')
        tf_fb2, _ = _init_textbox(i_slide, panel_x + Inches(0.3), panel_y + Inches(1.0),
                                   panel_w - Inches(0.6), Inches(2.0))
        run_fb2 = tf_fb2.paragraphs[0].add_run()
        _set_run(run_fb2, feedback_inc, font_name="Arial", size=11, color=INK)

        if correct_opt:
            _add_paragraph(tf_fb2, f"The correct answer is: {correct_opt['letter']}) {correct_opt['text']}",
                          font_name="Arial", size=10, bold=True, color=CHARTREUSE, space_before=12)

        _add_footer(i_slide, slide_idx + qi * 3 + 2, total_slides)
        slides.append(('incorrect', i_slide))

    return slides


def create_tips_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 9: Inclusive Practice Tip (s-tip) — chartreuse accents."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Inclusive practice")
    y = _add_slide_title(slide, screen.get('title', 'Inclusive Practice Tip'), y)

    # Tip card
    card_x = CONTENT_PAD_X
    card_y = Inches(2.6)
    card_w = Inches(12.0)
    card_h = Inches(3.8)

    # Chartreuse left border
    _add_rect(slide, card_x, card_y, Inches(0.06), card_h, fill_color=CHARTREUSE)
    # Card background
    _add_rect(slide, card_x + Inches(0.06), card_y, card_w - Inches(0.06), card_h,
              fill_color=CHARTREUSE_BG, border_color=NEUTRAL, border_width=0.5)

    # Sections: Do This / Why It Works / When to Use
    sections = ["Do This", "Why It Works", "When to Use"]
    inner_x = card_x + Inches(0.4)
    inner_y = card_y + Inches(0.3)
    inner_w = card_w - Inches(0.8)
    section_h = Inches(1.0)

    items = screen.get('text_items', [])
    items_clean = [i for i in items if not any(skip in i for skip in ['Audio:', 'Captions:', 'Image:', 'SME:'])]

    for i, sec_title in enumerate(sections):
        # Section title
        tf_title, _ = _init_textbox(slide, inner_x, inner_y, inner_w, Inches(0.3))
        run_t = tf_title.paragraphs[0].add_run()
        _set_run(run_t, sec_title, font_name="Arial Black", size=11, bold=True, color=CHARTREUSE)

        # Section content
        if i < len(items_clean):
            tf_content, _ = _init_textbox(slide, inner_x, inner_y + Inches(0.3), inner_w, Inches(0.6))
            run_c = tf_content.paragraphs[0].add_run()
            _set_run(run_c, items_clean[i], font_name="Arial", size=11, color=INK)

        inner_y += section_h + Inches(0.1)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_map_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 10: MAP Activity (s-map) — lilac accents, 3 input fields."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 My action planning")
    y = _add_slide_title(slide, screen.get('title', 'My Accessibility Plan (MAP)'), y)

    # Subtitle
    y = _add_lede(slide, "Reflect on what you have learned. This is private and not submitted.", y)

    # 3 input fields: Stop / Start / Continue
    field_labels = ["STOP", "START", "CONTINUE"]
    field_prompts = [
        "What will I stop doing that creates barriers?",
        "What will I start doing to be more inclusive?",
        "What will I continue doing that already works well?",
    ]
    field_x = CONTENT_PAD_X
    field_w = Inches(12.0)
    field_h = Inches(1.1)
    field_y = Inches(3.2)

    items = screen.get('text_items', [])
    items_clean = [i for i in items if not any(skip in i for skip in ['Audio:', 'Captions:', 'keyboard'])]

    for i, (label, prompt) in enumerate(zip(field_labels, field_prompts)):
        # Field background
        _add_rect(slide, field_x, field_y, field_w, field_h, fill_color=LILAC_BG, border_color=LILAC, border_width=0.5)
        # Lilac left accent
        _add_rect(slide, field_x, field_y, Inches(0.06), field_h, fill_color=LILAC)

        # Label
        tf_label, _ = _init_textbox(slide, field_x + Inches(0.25), field_y + Inches(0.1), Inches(2), Inches(0.3))
        run_l = tf_label.paragraphs[0].add_run()
        _set_run(run_l, label, font_name="Arial Black", size=10, bold=True, color=LILAC)

        # Prompt / text
        use_prompt = prompt
        if i < len(items_clean):
            use_prompt = items_clean[i] if items_clean[i] else prompt
        tf_prompt, _ = _init_textbox(slide, field_x + Inches(0.25), field_y + Inches(0.4),
                                      field_w - Inches(0.5), Inches(0.6))
        run_p = tf_prompt.paragraphs[0].add_run()
        _set_run(run_p, use_prompt, font_name="Arial", size=10, italic=True, color=MUTED)

        field_y += field_h + Inches(0.15)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_summary_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Slide 11: Key Takeaways (s-take) — 2x2 numbered takeaway cards."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Review")
    y = _add_slide_title(slide, "Key Takeaways", y)

    items = []
    for item in screen.get('text_items', []):
        clean = re.sub(r'^\d+\.\s*', '', item)
        if clean and 'Audio' not in clean and 'Captions' not in clean:
            items.append(clean)

    # 2x2 grid of takeaway cards
    card_w = Inches(5.5)
    card_h = Inches(1.5)
    start_x = CONTENT_PAD_X
    start_y = Inches(2.8)
    gap_x = Inches(0.5)
    gap_y = Inches(0.35)

    for i, item in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y_card = start_y + row * (card_h + gap_y)

        # Card
        _add_rect(slide, x, y_card, card_w, card_h, fill_color=WHITE, border_color=NEUTRAL, border_width=0.5)
        _add_rect(slide, x, y_card, card_w, Inches(0.04), fill_color=COBALT)

        # Number
        tf_num, _ = _init_textbox(slide, x + Inches(0.15), y_card + Inches(0.12), Inches(0.5), Inches(0.5))
        run_num = tf_num.paragraphs[0].add_run()
        _set_run(run_num, str(i + 1), font_name="Arial Black", size=28, bold=True, color=COBALT)

        # Text
        tf_item, _ = _init_textbox(slide, x + Inches(0.7), y_card + Inches(0.12),
                                    card_w - Inches(0.85), card_h - Inches(0.2))
        run_item = tf_item.paragraphs[0].add_run()
        _set_run(run_item, item, font_name="Arial", size=11, color=INK)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


def create_completion_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Final completion / resources slide — navy background with resources."""
    slide = _add_blank_slide(prs)

    # Navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Lilac accent line at top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), fill_color=LILAC)

    # UHN logo
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), CONTENT_PAD_X, Inches(0.5), height=Inches(0.5))

    # Completion title
    tf_title, _ = _init_textbox(slide, CONTENT_PAD_X, Inches(1.5), Inches(12), Inches(1.0))
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    subtitle = GUIDE_SUBTITLES.get(guide_num, "")
    _set_run(run, f"Guide {guide_num:02d} Complete", font_name="Arial Black", size=40, bold=True, color=WHITE)

    tf_sub, _ = _init_textbox(slide, CONTENT_PAD_X, Inches(2.6), Inches(12), Inches(0.5))
    p2 = tf_sub.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    _set_run(run2, subtitle, font_name="Arial", size=18, color=LILAC)

    # Resources list
    items = [i for i in screen.get('text_items', [])
             if 'Audio' not in i and 'Captions' not in i and 'badge' not in i.lower()]
    if items:
        res_y = Inches(3.5)
        tf_res_title, _ = _init_textbox(slide, Inches(3.0), res_y, Inches(7.333), Inches(0.4))
        p_rt = tf_res_title.paragraphs[0]
        p_rt.alignment = PP_ALIGN.CENTER
        run_rt = p_rt.add_run()
        _set_run(run_rt, "RESOURCES", font_name="Arial Black", size=10, bold=True, color=LILAC)

        res_y += Inches(0.5)
        for item in items[:6]:
            tf_r, _ = _init_textbox(slide, Inches(2.0), res_y, Inches(9.333), Inches(0.35))
            p_r = tf_r.paragraphs[0]
            p_r.alignment = PP_ALIGN.CENTER
            run_r = p_r.add_run()
            _set_run(run_r, item, font_name="Arial", size=11, color=RGBColor(0xCC, 0xCC, 0xDD))
            res_y += Inches(0.35)

    # "Next Guide" button
    btn_w = Inches(2.667)
    btn_x = (SLIDE_W - btn_w) / 2
    btn = _add_rounded_rect(slide, btn_x, Inches(6.0), btn_w, Inches(0.55), fill_color=RED)
    tf_btn, _ = _init_textbox(slide, btn_x, Inches(6.0), btn_w, Inches(0.55))
    tf_btn.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_btn = tf_btn.paragraphs[0]
    p_btn.alignment = PP_ALIGN.CENTER
    run_btn = p_btn.add_run()
    next_label = "NEXT GUIDE" if guide_num < 17 else "COMPLETE SERIES"
    _set_run(run_btn, next_label, font_name="Arial Black", size=12, bold=True, color=WHITE)

    # Bottom accent
    _add_rect(slide, Inches(0), FOOTER_Y + Inches(0.3), SLIDE_W, Inches(0.05), fill_color=LILAC)

    _add_notes(slide, screen)
    return slide


def create_reflection_slide(prs, screen, guide_num, slide_idx, total_slides):
    """Reflection slide — cobalt accents."""
    slide = _add_blank_slide(prs)
    section = _extract_section_name(screen)
    _add_topbar(slide, guide_num, section)

    section_num = screen.get('number', '').split('.')[0]
    y = _add_eyebrow(slide, f"Section {section_num} \u00b7 Reflection")
    y = _add_slide_title(slide, screen.get('title', 'Reflection'), y)

    items = screen.get('text_items', [])
    items_clean = [i for i in items if 'Audio' not in i and 'Captions' not in i and 'keyboard' not in i]

    content_y = Inches(2.8)
    for item in items_clean[:6]:
        # Reflection prompt with cobalt left border
        _add_rect(slide, CONTENT_PAD_X, content_y, Inches(0.06), Inches(0.6), fill_color=COBALT)
        _add_rect(slide, CONTENT_PAD_X + Inches(0.06), content_y, Inches(11.5), Inches(0.6),
                  fill_color=COBALT_BG)
        tf, _ = _init_textbox(slide, CONTENT_PAD_X + Inches(0.25), content_y + Inches(0.1),
                               Inches(11), Inches(0.4))
        run = tf.paragraphs[0].add_run()
        is_private = 'private' in item.lower() or 'not submitted' in item.lower()
        _set_run(run, item, font_name="Arial", size=12, italic=is_private, color=INK if not is_private else MUTED)
        content_y += Inches(0.75)

    _add_footer(slide, slide_idx, total_slides)
    _add_notes(slide, screen)
    return slide


# ---------------------------------------------------------------------------
# Main generation
# ---------------------------------------------------------------------------

def _count_total_slides(screens):
    """Pre-count total slides for progress dots."""
    count = 1  # title slide
    for screen in screens:
        stype = screen['type']
        if stype == 'scenario':
            count += 1 + len(screen.get('choices', []))
        elif stype == 'question':
            questions = screen.get('questions', [])
            count += max(len(questions), 1) * 3
        else:
            count += 1
    return count


def generate_ppt(guide_num: int, dry_run: bool = False):
    """Generate the full PPT from the master storyboard."""
    guide_folder = BUILD_OUTPUT / GUIDE_FOLDERS[guide_num]
    storyboard_path = guide_folder / "01-master-storyboard" / f"MASTER-STORYBOARD-GUIDE-{guide_num:02d}.md"

    if not storyboard_path.exists():
        print(f"Error: Storyboard not found: {storyboard_path}")
        sys.exit(1)

    print(f"Parsing storyboard: {storyboard_path.name}")
    data = parse_storyboard(storyboard_path)
    info = data['info']
    screens = data['screens']

    print(f"Found {len(screens)} screens")
    print(f"Course: {info.get('title', 'Unknown')}")
    print(f"Seat time: {info.get('seat_time', 'N/A')}")
    print()

    total_slides = _count_total_slides(screens)

    # Plan slides
    slide_plan = []
    slide_plan.append(('title', 'Course Title Slide', None))

    for screen in screens:
        stype = screen['type']
        label = f"Screen {screen['number']} -- {screen['title']}"

        if stype == 'scenario':
            slide_plan.append(('scenario_setup', f"{label} (Setup)", screen))
            for choice in screen.get('choices', []):
                slide_plan.append(('scenario_choice', f"  -> Choice {choice['letter']} ({choice['quality']})", screen))
        elif stype == 'question':
            questions = screen.get('questions', [])
            for qi, q in enumerate(questions):
                q_label = f"{label}"
                if len(questions) > 1:
                    q_label += f" Q{qi+1}"
                slide_plan.append(('question', f"{q_label} (Question)", screen))
                slide_plan.append(('correct', f"  -> Correct Feedback", screen))
                slide_plan.append(('incorrect', f"  -> Incorrect Feedback", screen))
        else:
            slide_plan.append((stype, label, screen))

    print(f"Slide plan ({len(slide_plan)} slides):")
    print("-" * 60)
    for i, (stype, label, _) in enumerate(slide_plan, 1):
        print(f"  {i:3d}. [{stype:16s}] {label}")
    print("-" * 60)
    print()

    if dry_run:
        print("Dry run -- no PPT generated.")
        return

    # Generate PPT
    print("Generating PPT (mockup-matched, shapes-from-scratch)...")

    # Use template to get slide masters, but build on blank layouts
    if TEMPLATE_PATH.exists():
        prs = Presentation(str(TEMPLATE_PATH))
        # Remove any existing slides from the template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    else:
        print("Warning: UHN template not found, using blank presentation")
        prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_idx = 0

    # 1. Title slide
    create_title_slide(prs, info, guide_num, total_slides)
    print("  [ok] Title slide")
    slide_idx += 1

    # 2. Screen slides
    for screen in screens:
        stype = screen['type']
        label = f"Screen {screen['number']} -- {screen['title']}"

        if stype == 'welcome':
            create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (welcome)")
            slide_idx += 1

        elif stype == 'objectives':
            create_objectives_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (objectives)")
            slide_idx += 1

        elif stype == 'scenario':
            scenario_slides = create_scenario_slides(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (scenario: {len(scenario_slides)} slides)")
            slide_idx += len(scenario_slides)

        elif stype == 'question':
            q_slides = create_question_slides(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (questions: {len(q_slides)} slides)")
            slide_idx += len(q_slides)

        elif stype == 'action_planning':
            create_map_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (MAP)")
            slide_idx += 1

        elif stype == 'reflection':
            create_reflection_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (reflection)")
            slide_idx += 1

        elif stype == 'tips':
            create_tips_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (tip)")
            slide_idx += 1

        elif stype == 'summary':
            create_summary_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (takeaways)")
            slide_idx += 1

        elif stype == 'completion':
            create_completion_slide(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (completion)")
            slide_idx += 1

        else:
            create_content_slide_generic(prs, screen, guide_num, slide_idx, total_slides)
            print(f"  [ok] {label} (content)")
            slide_idx += 1

    # Save
    output_dir = guide_folder / "02-ppt-storyboard"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"UHN_Guide{guide_num:02d}_Storyboard.pptx"
    prs.save(str(output_path))

    actual_slides = len(prs.slides)
    print(f"\n[ok] PPT saved: {output_path}")
    print(f"  Total slides: {actual_slides}")
    print(f"  Dimensions: 13.333\" x 7.5\" (16:9 widescreen)")
    print(f"  Size: {output_path.stat().st_size / 1024:.0f} KB")


def main():
    parser = argparse.ArgumentParser(description="Generate PPT storyboard for UHN Accessibility First courses")
    parser.add_argument("--guide", type=int, required=True, help="Guide number (1-18)")
    parser.add_argument("--dry-run", action="store_true", help="Preview slide plan without generating")
    args = parser.parse_args()

    if args.guide not in GUIDE_FOLDERS:
        print(f"Error: Guide {args.guide} not found. Valid: 1-18")
        sys.exit(1)

    generate_ppt(args.guide, dry_run=args.dry_run)


if __name__ == "__main__":
    main()
